from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys
import traceback

file_path = r'C:\Users\Windows\Desktop\PLANIFICACION 2026\PDC_5TO DE PRIMARIA\PRIMER TRIMESTRE_2026\SOCIALES _ 2026\Diapositivas\Tema 1 - El estudio de la historia.pptx'

try:
    prs = Presentation(file_path)
except Exception as e:
    print(f"Error loading: {e}")
    sys.exit(1)

print(f"Total Slides: {len(prs.slides)}")
print(f"Slide dimensions: {prs.slide_width / 914400:.2f} x {prs.slide_height / 914400:.2f} inches")

for i, slide in enumerate(list(prs.slides)[:4]):
    print(f"\n--- SLIDE {i} ---")
    try:
        if hasattr(slide, 'background') and slide.background and hasattr(slide.background, 'fill'):
            fill = slide.background.fill
            if fill.type == 1 and hasattr(fill.fore_color, 'rgb'):
                print(f"Background: #{fill.fore_color.rgb}")
    except Exception:
        pass

    for j, shape in enumerate(slide.shapes):
        t = shape.shape_type
        type_str = str(t).split(' ')[0]
        l = getattr(shape, 'left', 0) / 914400 if getattr(shape, 'left', 0) else 0
        w = getattr(shape, 'width', 0) / 914400 if getattr(shape, 'width', 0) else 0
        h = getattr(shape, 'height', 0) / 914400 if getattr(shape, 'height', 0) else 0
        print(f"  Shape {j} ({type_str}) - w:{w:.2f} h:{h:.2f} l:{l:.2f}")

        if shape.has_text_frame:
            for p_idx, p in enumerate(shape.text_frame.paragraphs):
                if not p.text.strip(): continue
                print(f"    P{p_idx}: {repr(p.text[:60])}")
                try:
                    if p.runs:
                        font = p.runs[0].font
                        f_name = font.name
                        f_size = font.size.pt if getattr(font.size, 'pt', None) else "N/A"
                        f_color = "Theme/Default"
                        if hasattr(font.color, 'rgb') and font.color.rgb:
                            f_color = f"#{font.color.rgb}"
                        print(f"      Style: Font={f_name}, Size={f_size}, Color={f_color}")
                except Exception:
                    pass
