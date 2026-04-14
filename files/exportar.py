"""
exportar.py — Módulo de exportación PRIA v5.4
Genera archivos DOCX y PPTX desde los datos de sesión.

Uso desde app_ui.py:
    from exportar import exportar_docx, exportar_pptx, render_panel_exportacion
"""

import io
import json
from datetime import datetime

# ── python-docx ──────────────────────────────────────────────────────────────
from docx import Document
from docx.shared import Pt, RGBColor, Inches, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

# ── python-pptx ──────────────────────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt, Emu
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as Inch

# ─────────────────────────────────────────────────────────────────────────────
# PALETA DE COLORES PRIA
# ─────────────────────────────────────────────────────────────────────────────
VERDE_OSCURO  = RGBColor(0x1B, 0x5E, 0x20)   # #1B5E20
VERDE_MEDIO   = RGBColor(0x2E, 0x7D, 0x32)   # #2E7D32
VERDE_CLARO   = RGBColor(0xC8, 0xE6, 0xC9)   # #C8E6C9
VERDE_PALIDO  = RGBColor(0xF1, 0xF8, 0xE9)   # #F1F8E9
BLANCO        = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_TEXTO    = RGBColor(0x33, 0x33, 0x33)

# pptx equivalents
P_VERDE_OSCURO = PRGBColor(0x1B, 0x5E, 0x20)
P_VERDE_MEDIO  = PRGBColor(0x2E, 0x7D, 0x32)
P_VERDE_CLARO  = PRGBColor(0xC8, 0xE6, 0xC9)
P_VERDE_PALIDO = PRGBColor(0xF1, 0xF8, 0xE9)
P_BLANCO       = PRGBColor(0xFF, 0xFF, 0xFF)
P_GRIS         = PRGBColor(0x33, 0x33, 0x33)

FONT_TITULO = "Calibri"
FONT_CUERPO = "Calibri"

# ─────────────────────────────────────────────────────────────────────────────
# HELPERS DOCX
# ─────────────────────────────────────────────────────────────────────────────
def _set_cell_bg(cell, hex_color: str):
    """Pinta el fondo de una celda de tabla en DOCX."""
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd  = OxmlElement("w:shd")
    shd.set(qn("w:val"),   "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"),  hex_color)
    tcPr.append(shd)


def _heading(doc: Document, text: str, level: int = 1):
    """Añade un título con estilo PRIA."""
    p    = doc.add_heading(text, level=level)
    run  = p.runs[0] if p.runs else p.add_run(text)
    run.font.color.rgb = VERDE_OSCURO
    run.font.name      = FONT_TITULO
    run.font.bold      = True
    run.font.size      = Pt({1: 18, 2: 14, 3: 12}.get(level, 11))
    p.paragraph_format.space_before = Pt(12)
    p.paragraph_format.space_after  = Pt(4)
    return p


def _parrafo(doc: Document, text: str, bold: bool = False, italic: bool = False,
             color: RGBColor = None, size: int = 11):
    """Añade un párrafo normal con estilo PRIA."""
    p   = doc.add_paragraph()
    run = p.add_run(text)
    run.font.name   = FONT_CUERPO
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color or GRIS_TEXTO
    p.paragraph_format.space_after = Pt(4)
    return p


def _bullet(doc: Document, text: str, level: int = 0):
    """Añade un ítem de lista con viñeta."""
    p   = doc.add_paragraph(style="List Bullet")
    run = p.add_run(text)
    run.font.name  = FONT_CUERPO
    run.font.size  = Pt(11)
    run.font.color.rgb = GRIS_TEXTO
    p.paragraph_format.left_indent = Inches(0.25 * (level + 1))
    return p


def _tabla_simple(doc: Document, filas: list[dict], color_header: str = "1B5E20"):
    """Crea una tabla con encabezados desde una lista de dicts."""
    if not filas:
        return
    columnas = list(filas[0].keys())
    tabla    = doc.add_table(rows=1, cols=len(columnas))
    tabla.style = "Table Grid"
    tabla.alignment = WD_TABLE_ALIGNMENT.LEFT

    # Encabezados
    hdr_cells = tabla.rows[0].cells
    for i, col in enumerate(columnas):
        hdr_cells[i].text = col
        _set_cell_bg(hdr_cells[i], color_header)
        run = hdr_cells[i].paragraphs[0].runs[0]
        run.font.color.rgb = BLANCO
        run.font.bold      = True
        run.font.name      = FONT_TITULO
        run.font.size      = Pt(10)

    # Datos
    for fila in filas:
        row_cells = tabla.add_row().cells
        for i, col in enumerate(columnas):
            row_cells[i].text = str(fila.get(col, ""))
            _set_cell_bg(row_cells[i], "F1F8E9")
            run = row_cells[i].paragraphs[0].runs[0]
            run.font.color.rgb = VERDE_OSCURO
            run.font.name      = FONT_CUERPO
            run.font.size      = Pt(10)

    doc.add_paragraph()  # espacio tras tabla


def _separador(doc: Document):
    p = doc.add_paragraph()
    pPr = p._p.get_or_add_pPr()
    pBdr = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"),   "single")
    bottom.set(qn("w:sz"),    "6")
    bottom.set(qn("w:space"), "1")
    bottom.set(qn("w:color"), "2E7D32")
    pBdr.append(bottom)
    pPr.append(pBdr)


def _pie_pagina(doc: Document, titulo: str):
    """Añade pie de página con título y fecha."""
    section  = doc.sections[0]
    footer   = section.footer
    p        = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
    p.clear()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(f"PRIA — Método Palma-Ribera  |  {titulo}  |  {datetime.now().strftime('%d/%m/%Y')}")
    run.font.name      = FONT_CUERPO
    run.font.size      = Pt(8)
    run.font.color.rgb = VERDE_MEDIO
    run.font.italic    = True


def _encabezado_doc(doc: Document, titulo: str, subtitulo: str = "", tema: str = ""):
    """Encabezado visual de portada dentro del documento."""
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("🦉 PRIA — Método Palma-Ribera")
    run.font.name      = FONT_TITULO
    run.font.size      = Pt(10)
    run.font.color.rgb = VERDE_MEDIO
    run.font.italic    = True

    h = doc.add_heading(titulo, level=1)
    for run in h.runs:
        run.font.color.rgb = VERDE_OSCURO
        run.font.name      = FONT_TITULO
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER

    if subtitulo:
        ps = doc.add_paragraph()
        ps.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = ps.add_run(subtitulo)
        r.font.name      = FONT_CUERPO
        r.font.size      = Pt(12)
        r.font.color.rgb = VERDE_MEDIO

    if tema:
        pt = doc.add_paragraph()
        pt.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = pt.add_run(f"Tema: {tema}")
        r.font.name      = FONT_CUERPO
        r.font.size      = Pt(11)
        r.font.bold      = True
        r.font.color.rgb = GRIS_TEXTO

    _separador(doc)
    doc.add_paragraph()


def _markdown_a_docx(doc: Document, texto: str):
    """
    Convierte texto Markdown básico a párrafos DOCX.
    Soporta: ##, ###, **, *, - , | (tablas simples), >
    """
    if not texto:
        return
    lines = texto.split("\n")
    tabla_buf = []

    def flush_tabla():
        if not tabla_buf:
            return
        # Filtrar líneas separadoras (---|---)
        filas_validas = [l for l in tabla_buf if not all(c in "-| " for c in l)]
        if not filas_validas:
            tabla_buf.clear()
            return
        cols = [c.strip() for c in filas_validas[0].split("|") if c.strip()]
        rows_data = []
        for row_line in filas_validas[1:]:
            vals = [c.strip() for c in row_line.split("|") if c.strip()]
            if vals:
                rows_data.append(dict(zip(cols, vals + [""] * (len(cols) - len(vals)))))
        if rows_data:
            _tabla_simple(doc, rows_data)
        tabla_buf.clear()

    for line in lines:
        stripped = line.strip()

        # Tabla markdown
        if stripped.startswith("|"):
            tabla_buf.append(stripped)
            continue
        else:
            flush_tabla()

        if not stripped:
            doc.add_paragraph()
            continue

        if stripped.startswith("### "):
            _heading(doc, stripped[4:], level=3)
        elif stripped.startswith("## "):
            _heading(doc, stripped[3:], level=2)
        elif stripped.startswith("# "):
            _heading(doc, stripped[2:], level=1)
        elif stripped.startswith("> "):
            p   = doc.add_paragraph()
            run = p.add_run(stripped[2:])
            run.font.italic    = True
            run.font.color.rgb = VERDE_MEDIO
            run.font.size      = Pt(11)
            p.paragraph_format.left_indent = Inches(0.4)
        elif stripped.startswith("- ") or stripped.startswith("* "):
            _bullet(doc, stripped[2:])
        elif stripped.startswith("**") and stripped.endswith("**"):
            _parrafo(doc, stripped.strip("*"), bold=True)
        else:
            # Inline bold/italic dentro del texto
            p   = doc.add_paragraph()
            run = p.add_run(stripped.replace("**","").replace("__",""))
            run.font.name      = FONT_CUERPO
            run.font.size      = Pt(11)
            run.font.color.rgb = GRIS_TEXTO

    flush_tabla()


# ─────────────────────────────────────────────────────────────────────────────
# GENERADORES DOCX — uno por tipo de documento
# ─────────────────────────────────────────────────────────────────────────────
def _nuevo_doc() -> Document:
    doc = Document()
    # Márgenes 1 pulgada
    for section in doc.sections:
        section.top_margin    = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin   = Inches(1.1)
        section.right_margin  = Inches(1.1)
    # Fuente por defecto
    style = doc.styles["Normal"]
    style.font.name = FONT_CUERPO
    style.font.size = Pt(11)
    return doc


def generar_docx_sintesis(res_m0a: dict, diagnosticos: str) -> bytes:
    """DOCX — Síntesis Curricular de Unidad (M0a)."""
    doc = _nuevo_doc()
    unidad = res_m0a.get("unidad_sintetizada", {})
    titulo = unidad.get("titulo", "Síntesis de Unidad")

    _encabezado_doc(doc, "Síntesis Curricular de Unidad", titulo)
    _pie_pagina(doc, "Síntesis Curricular")

    # Perfil del aula
    _heading(doc, "🩺 Perfil Neuro-Inclusivo del Aula", 2)
    _parrafo(doc, diagnosticos)
    doc.add_paragraph()

    # Temas desarrollados
    _heading(doc, "📑 Temas de la Unidad", 2)
    for tema in unidad.get("temas_desarrollados", []):
        _heading(doc, f"📍 {tema.get('nombre','Tema')}", 3)

        _parrafo(doc, "Conceptos Clave:", bold=True, color=VERDE_OSCURO)
        for c in tema.get("conceptos_clave") or tema.get("conceptos") or []:
            _bullet(doc, str(c))

        _parrafo(doc, "Inteligencias Múltiples sugeridas:", bold=True, color=VERDE_OSCURO)
        for i in tema.get("inteligencias_sugeridas") or tema.get("inteligencias") or []:
            _bullet(doc, str(i))
        doc.add_paragraph()

    _separador(doc)

    # Notas DUA
    _heading(doc, "👨‍🏫 Notas DUA del Docente", 2)
    _parrafo(doc, unidad.get("notas_docente", "—"))
    doc.add_paragraph()

    # Proyecto ABP
    _heading(doc, "🚀 Proyecto ABP — Resumen", 2)
    _parrafo(doc, unidad.get("proyecto_pbl", "—"), italic=True)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def generar_docx_abp(res_m0b: str, res_m0c: str, titulo_unidad: str) -> bytes:
    """DOCX — Proyecto ABP completo (M0b + M0c)."""
    doc = _nuevo_doc()
    _encabezado_doc(doc, "Proyecto ABP Completo", titulo_unidad)
    _pie_pagina(doc, "Proyecto ABP")

    _heading(doc, "📋 Diseño del Proyecto ABP", 2)
    _markdown_a_docx(doc, res_m0b or "No generado.")
    _separador(doc)
    doc.add_paragraph()

    _heading(doc, "📊 Rúbrica y Fichas de Proceso", 2)
    _markdown_a_docx(doc, res_m0c or "No generado.")

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def generar_docx_plan_clase(res_m1a: dict, tema: str, diagnosticos: str) -> bytes:
    """DOCX — Plan de Clase de 45 minutos (M1a)."""
    doc = _nuevo_doc()
    _encabezado_doc(doc, "Plan de Clase", "Secuencia Didáctica de 45 minutos", tema)
    _pie_pagina(doc, "Plan de Clase")

    # Verbos Bloom
    verbos = res_m1a.get("mapa_cognitivo", {}).get("verbos", [])
    if verbos:
        _heading(doc, "🎯 Verbos Operativos (Taxonomía Bloom)", 2)
        for v in verbos:
            _bullet(doc, str(v))
        doc.add_paragraph()

    # Inteligencias Múltiples
    im = res_m1a.get("inteligencias_multiples", [])
    if im:
        _heading(doc, "🧩 Inteligencias Múltiples", 2)
        _tabla_simple(doc, im)

    # Secuencia didáctica
    _heading(doc, "⏱️ Secuencia Didáctica", 2)
    for bloque in res_m1a.get("secuencia_didactica", {}).get("bloques", []):
        nombre = bloque.get("nombre", "Bloque")
        dur    = bloque.get("duracion", "?")
        _heading(doc, f"{nombre} — {dur} minutos", 3)
        _parrafo(doc, bloque.get("objetivo", ""))
        if "nota" in bloque:
            p   = doc.add_paragraph()
            run = p.add_run(f"💡 {bloque['nota']}")
            run.font.italic    = True
            run.font.color.rgb = VERDE_MEDIO
            run.font.size      = Pt(10)
        doc.add_paragraph()

    _separador(doc)

    # DUA
    _heading(doc, "🛡️ Directrices DUA — Neuro-Inclusión", 2)
    for d in res_m1a.get("dua_neuroinclusion", []):
        _bullet(doc, str(d))
    doc.add_paragraph()

    # Adaptaciones
    adapt = res_m1a.get("tabla_adaptaciones_clase", [])
    if adapt:
        _heading(doc, "♿ Adaptaciones por Diagnóstico", 2)
        _tabla_simple(doc, adapt)

    # Perfil del aula
    _heading(doc, "👥 Perfil del Aula", 2)
    _parrafo(doc, res_m1a.get("perfil_aula_resumido", diagnosticos))

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def generar_docx_ficha(res_m1c: dict, tema: str) -> bytes:
    """DOCX — Ficha Gamificada (M1c)."""
    doc = _nuevo_doc()
    raw   = res_m1c
    ficha = raw.get("ficha_trabajo", raw) if isinstance(raw, dict) else {}

    _encabezado_doc(doc, ficha.get("titulo_gancho", "Ficha Gamificada"), tema)
    _pie_pagina(doc, "Ficha Gamificada")

    _parrafo(doc, ficha.get("historia_gancho", ""), italic=True)
    _separador(doc)
    doc.add_paragraph()

    misiones = ficha.get("misiones", {})

    # Oráculo
    _heading(doc, "🔮 Misión 1: El Oráculo", 2)
    for q in misiones.get("oraculo", []):
        _parrafo(doc, q.get("pregunta", ""), bold=True)
        for op in q.get("opciones", []):
            _bullet(doc, str(op))
        doc.add_paragraph()

    # Puente
    _heading(doc, "🌉 Misión 2: El Puente", 2)
    pares = [{"Palabra": p.get("palabra",""), "Significado": p.get("significado","")}
             for p in misiones.get("puente", [])]
    if pares:
        _tabla_simple(doc, pares)

    # Sopa
    _heading(doc, "🥣 Misión 3: La Sopa", 2)
    for p in misiones.get("sopa", []):
        _bullet(doc, str(p))
    doc.add_paragraph()

    # Pergamino
    _heading(doc, "📜 Misión 4: El Pergamino", 2)
    per = misiones.get("pergamino", {})
    _parrafo(doc, per.get("frase_con_espacios", ""), italic=True)
    _parrafo(doc, f"Palabras secretas: {', '.join(per.get('palabras_secretas', []))}", bold=True)
    doc.add_paragraph()

    # Lienzo
    _heading(doc, "🎨 Misión 5: El Lienzo", 2)
    _parrafo(doc, misiones.get("lienzo", ""))
    doc.add_paragraph()

    # Adaptaciones
    adapt = ficha.get("adaptaciones_por_mision", [])
    if adapt:
        _separador(doc)
        _heading(doc, "♿ Adaptaciones por Misión", 2)
        _tabla_simple(doc, adapt)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def generar_docx_evaluaciones(res_m2a: str, res_m2b: str, tema: str) -> bytes:
    """DOCX — Pop Quiz DUA + Guía del Tutor (M2a + M2b)."""
    doc = _nuevo_doc()
    _encabezado_doc(doc, "Instrumentos de Evaluación", "Pop Quiz DUA + Guía del Tutor", tema)
    _pie_pagina(doc, "Evaluaciones")

    _heading(doc, "🎲 Pop Quiz DUA", 2)
    _markdown_a_docx(doc, res_m2a or "No generado.")
    _separador(doc)
    doc.add_paragraph()

    _heading(doc, "📊 Panel de Control del Tutor", 2)
    _markdown_a_docx(doc, res_m2b or "No generado.")

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ─────────────────────────────────────────────────────────────────────────────
# GENERADOR PPTX — Presentación de clase (M1b)
# ─────────────────────────────────────────────────────────────────────────────
def generar_pptx_diapositivas(res_m1b, tema: str, unidad: str) -> bytes:
    """
    PPTX — Presentación de clase desde la estructura JSON de M1b.
    Diseño Forest & Moss: verde oscuro + crema, tipografía Calibri.
    """
    prs = Presentation()
    prs.slide_width  = PInches(13.33)
    prs.slide_height = PInches(7.5)

    # Layouts: 0=título+contenido, 6=en blanco
    LAYOUT_TITULO   = prs.slide_layouts[0]
    LAYOUT_CONTENIDO = prs.slide_layouts[1]
    LAYOUT_BLANCO   = prs.slide_layouts[6]

    # Normalizar estructura de slides
    if isinstance(res_m1b, list):
        slides_data = res_m1b
    elif isinstance(res_m1b, dict):
        slides_data = (res_m1b.get("diapositivas")
                       or res_m1b.get("slides")
                       or (list(res_m1b.values())[0] if res_m1b else []))
    else:
        slides_data = []

    def _bg_verde_oscuro(slide):
        """Rellena el fondo de la diapositiva con verde oscuro."""
        bg   = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = P_VERDE_OSCURO

    def _bg_crema(slide):
        """Fondo crema/blanco para diapositivas de contenido."""
        bg   = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = PRGBColor(0xF9, 0xF7, 0xF0)

    def _add_textbox(slide, left, top, width, height, text,
                     font_size=18, bold=False, color=None,
                     align=PP_ALIGN.LEFT, italic=False):
        txBox = slide.shapes.add_textbox(
            PInches(left), PInches(top), PInches(width), PInches(height)
        )
        tf  = txBox.text_frame
        tf.word_wrap = True
        p   = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text           = text
        run.font.size      = PPt(font_size)
        run.font.bold      = bold
        run.font.italic    = italic
        run.font.color.rgb = color or P_BLANCO
        return txBox

    def _franja_inferior(slide, texto: str):
        """Franja verde medio en la parte inferior con texto de nota."""
        franja = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            PInches(0), PInches(6.6),
            PInches(13.33), PInches(0.9)
        )
        franja.fill.solid()
        franja.fill.fore_color.rgb = P_VERDE_MEDIO
        franja.line.fill.background()
        tf  = franja.text_frame
        tf.word_wrap = True
        p   = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text           = texto[:200]
        run.font.size      = PPt(11)
        run.font.color.rgb = P_BLANCO
        run.font.italic    = True

    # ── Slide de portada del tema ────────────────────────────────────────────
    portada = prs.slides.add_slide(LAYOUT_BLANCO)
    _bg_verde_oscuro(portada)

    # Franja decorativa lateral izquierda
    barra = portada.shapes.add_shape(
        1, PInches(0), PInches(0), PInches(0.4), PInches(7.5)
    )
    barra.fill.solid()
    barra.fill.fore_color.rgb = P_VERDE_CLARO
    barra.line.fill.background()

    _add_textbox(portada, 0.7, 0.5, 12, 0.5,
                 "🦉 PRIA — Método Palma-Ribera",
                 font_size=12, italic=True, color=P_VERDE_CLARO)
    _add_textbox(portada, 0.7, 1.2, 11, 2.0,
                 unidad, font_size=22, bold=False, color=P_VERDE_CLARO)
    _add_textbox(portada, 0.7, 2.8, 11, 2.0,
                 tema,   font_size=38, bold=True,  color=P_BLANCO)
    _add_textbox(portada, 0.7, 5.8, 11, 0.5,
                 datetime.now().strftime("%d de %B de %Y"),
                 font_size=11, color=P_VERDE_CLARO, italic=True)

    # ── Diapositivas de contenido ────────────────────────────────────────────
    for i, slide_data in enumerate(slides_data):
        if not isinstance(slide_data, dict):
            continue

        titulo_slide = (slide_data.get("titulo") or slide_data.get("title")
                        or f"Diapositiva {i+1}")
        guion        = (slide_data.get("guion_docente") or slide_data.get("guion") or "")
        texto_pan    = (slide_data.get("texto_pantalla") or slide_data.get("texto") or "")
        prompt_img   = (slide_data.get("prompt_imagen") or slide_data.get("prompt") or "")
        es_pausa     = "kinest" in titulo_slide.lower() or "pausa" in titulo_slide.lower()

        slide = prs.slides.add_slide(LAYOUT_BLANCO)

        if es_pausa:
            # Diapositiva de pausa kinestésica — fondo verde claro
            bg = slide.background; bg.fill.solid()
            bg.fill.fore_color.rgb = P_VERDE_CLARO
            _add_textbox(slide, 1, 2.5, 11, 1.5,
                         "🏃 " + titulo_slide,
                         font_size=36, bold=True, color=P_VERDE_OSCURO,
                         align=PP_ALIGN.CENTER)
            _add_textbox(slide, 1, 4.2, 11, 1.5,
                         texto_pan, font_size=16, color=P_VERDE_OSCURO,
                         align=PP_ALIGN.CENTER)
        else:
            _bg_crema(slide)

            # Barra de título
            barra_t = slide.shapes.add_shape(
                1, PInches(0), PInches(0), PInches(13.33), PInches(1.2)
            )
            barra_t.fill.solid()
            barra_t.fill.fore_color.rgb = P_VERDE_OSCURO
            barra_t.line.fill.background()

            # Número de slide
            _add_textbox(slide, 12.5, 0.1, 0.7, 0.4,
                         str(i + 1), font_size=11, color=P_VERDE_CLARO)

            # Título
            _add_textbox(slide, 0.3, 0.15, 12, 0.9,
                         titulo_slide, font_size=22, bold=True, color=P_BLANCO)

            # Texto de pantalla (columna izquierda)
            _add_textbox(slide, 0.3, 1.4, 7.5, 4.8,
                         texto_pan, font_size=15, color=P_GRIS)

            # Prompt de imagen (columna derecha — como referencia visual)
            if prompt_img:
                caja_img = slide.shapes.add_shape(
                    1, PInches(8.1), PInches(1.4), PInches(4.8), PInches(4.5)
                )
                caja_img.fill.solid()
                caja_img.fill.fore_color.rgb = P_VERDE_PALIDO
                caja_img.line.color.rgb      = P_VERDE_CLARO

                tf  = caja_img.text_frame
                tf.word_wrap = True
                p   = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT

                r1 = p.add_run()
                r1.text           = "🖼 Imagen sugerida:\n"
                r1.font.size      = PPt(10)
                r1.font.bold      = True
                r1.font.color.rgb = P_VERDE_OSCURO

                r2 = p.add_run()
                r2.text           = prompt_img[:250]
                r2.font.size      = PPt(9)
                r2.font.italic    = True
                r2.font.color.rgb = P_VERDE_MEDIO

            # Franja inferior con guion del docente
            if guion:
                _franja_inferior(slide, f"📢 Docente: {guion[:180]}")

    # ── Slide de cierre ──────────────────────────────────────────────────────
    cierre = prs.slides.add_slide(LAYOUT_BLANCO)
    _bg_verde_oscuro(cierre)
    barra_c = cierre.shapes.add_shape(
        1, PInches(0), PInches(0), PInches(0.4), PInches(7.5)
    )
    barra_c.fill.solid()
    barra_c.fill.fore_color.rgb = P_VERDE_CLARO
    barra_c.line.fill.background()

    _add_textbox(cierre, 0.7, 2.5, 12, 1.5,
                 "¡Excelente trabajo hoy!", font_size=36, bold=True,
                 color=P_BLANCO, align=PP_ALIGN.CENTER)
    _add_textbox(cierre, 0.7, 4.2, 12, 0.8,
                 "🦉 PRIA — Método Palma-Ribera",
                 font_size=14, color=P_VERDE_CLARO,
                 align=PP_ALIGN.CENTER, italic=True)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ─────────────────────────────────────────────────────────────────────────────
# PANEL DE EXPORTACIÓN — para insertar en app_ui.py con st.expander
# ─────────────────────────────────────────────────────────────────────────────
def render_panel_exportacion(ss, diagnosticos: str):
    """
    Renderiza el panel completo de exportación en Streamlit.
    Llama con: render_panel_exportacion(st.session_state, DIAGNOSTICOS)
    """
    import streamlit as st

    unidad_titulo = ""
    if ss.res_m0a and "unidad_sintetizada" in ss.res_m0a:
        unidad_titulo = ss.res_m0a["unidad_sintetizada"].get("titulo", "Unidad")

    tema_activo = ss.get("tema_activo", "tema")
    fecha_str   = datetime.now().strftime("%Y%m%d")

    st.markdown("### 📤 Exportar Documentos")
    st.caption("Solo se habilitan los documentos que ya han sido generados en esta sesión.")

    # ── Fila 1: Documentos de Unidad ─────────────────────────────────────────
    st.markdown("#### 🏗️ Plan de Unidad y ABP")
    c1, c2, c3 = st.columns(3)

    with c1:
        if ss.res_m0a:
            data = generar_docx_sintesis(ss.res_m0a, diagnosticos)
            st.download_button(
                "📄 Síntesis Curricular (.docx)",
                data=data,
                file_name=f"PRIA_Sintesis_{fecha_str}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True
            )
        else:
            st.button("📄 Síntesis Curricular (.docx)", disabled=True, use_container_width=True,
                      help="Genera la Síntesis Curricular primero.")

    with c2:
        if ss.res_m0b:
            data = generar_docx_abp(ss.res_m0b, ss.res_m0c or "", unidad_titulo)
            st.download_button(
                "📋 Proyecto ABP completo (.docx)",
                data=data,
                file_name=f"PRIA_ABP_{fecha_str}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True
            )
        else:
            st.button("📋 Proyecto ABP completo (.docx)", disabled=True, use_container_width=True,
                      help="Diseña el Proyecto ABP completo primero.")

    with c3:
        st.empty()

    # ── Fila 2: Documentos de Clase ───────────────────────────────────────────
    st.markdown("#### 🚀 Plan de Clase y Fichas")
    c4, c5, c6 = st.columns(3)

    with c4:
        if ss.res_m1a:
            data = generar_docx_plan_clase(ss.res_m1a, tema_activo, diagnosticos)
            st.download_button(
                "📄 Plan de Clase (.docx)",
                data=data,
                file_name=f"PRIA_PlanClase_{fecha_str}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True
            )
        else:
            st.button("📄 Plan de Clase (.docx)", disabled=True, use_container_width=True,
                      help="Genera el Plan de Clase primero.")

    with c5:
        if ss.res_m1b:
            data = generar_pptx_diapositivas(ss.res_m1b, tema_activo, unidad_titulo)
            st.download_button(
                "🖼️ Diapositivas (.pptx)",
                data=data,
                file_name=f"PRIA_Slides_{fecha_str}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )
        else:
            st.button("🖼️ Diapositivas (.pptx)", disabled=True, use_container_width=True,
                      help="Genera las Diapositivas primero.")

    with c6:
        if ss.res_m1c:
            data = generar_docx_ficha(ss.res_m1c, tema_activo)
            st.download_button(
                "🎮 Ficha Gamificada (.docx)",
                data=data,
                file_name=f"PRIA_FichaGamificada_{fecha_str}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True
            )
        else:
            st.button("🎮 Ficha Gamificada (.docx)", disabled=True, use_container_width=True,
                      help="Genera la Ficha Gamificada primero.")

    # ── Fila 3: Evaluaciones ──────────────────────────────────────────────────
    st.markdown("#### 📝 Evaluaciones")
    c7, c8, c9 = st.columns(3)

    with c7:
        if ss.res_m2a:
            data = generar_docx_evaluaciones(ss.res_m2a, ss.res_m2b or "", tema_activo)
            st.download_button(
                "📝 Pop Quiz + Guía Tutor (.docx)",
                data=data,
                file_name=f"PRIA_Evaluaciones_{fecha_str}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True
            )
        else:
            st.button("📝 Pop Quiz + Guía Tutor (.docx)", disabled=True, use_container_width=True,
                      help="Genera el Pop Quiz primero.")

    with c8:
        st.empty()
    with c9:
        st.empty()
