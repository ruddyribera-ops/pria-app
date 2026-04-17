"""
exporters/pptx.py - PowerPoint Export Functions
================================================
PPTX helpers and generar_pptx_diapositivas function.
"""

import io
import os
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches as PIn, Pt as PPt
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.text import PP_ALIGN


# ─── LOGO PATH ────────────────────────────────────────────────────────────────
_DIR = Path(__file__).parent.parent
LOGO_PATH = str(_DIR / "logo_laspalmas.png")
if not os.path.exists(LOGO_PATH):
    LOGO_PATH = str(Path.cwd() / "logo_laspalmas.png")


# ─── COLORES ─────────────────────────────────────────────────────────────────
P_VERDE = PRGB(0x1B, 0x5E, 0x20)
P_NEGRO = PRGB(0x1A, 0x1A, 0x1A)
P_BLANCO = PRGB(0xFF, 0xFF, 0xFF)
P_GRIS = PRGB(0x55, 0x55, 0x55)
P_AMARILLO = PRGB(0xFF, 0xEE, 0x00)

FONT_TITLE = "Impact"
FONT_BODY = "Calibri"

# ─── SLIDE DIMENSIONS ─────────────────────────────────────────────────────────
W, H = 10.0, 5.625

_MESES_ES = {
    1: "enero",
    2: "febrero",
    3: "marzo",
    4: "abril",
    5: "mayo",
    6: "junio",
    7: "julio",
    8: "agosto",
    9: "septiembre",
    10: "octubre",
    11: "noviembre",
    12: "diciembre",
}


def _fecha_es() -> str:
    from datetime import datetime

    d = datetime.now()
    return f"{d.day:02d} de {_MESES_ES[d.month]} de {d.year}"


def _extract_scene(prompt: str) -> str:
    """Extract only PARTE 2 (the scene) from the three-part M1b image prompt."""
    p = str(prompt).strip()
    for marker in ("Square format.", "square format."):
        if marker in p:
            p = p[p.index(marker) + len(marker) :].strip().lstrip(".")
            break
    for marker in ("Consistent illustration", "consistent illustration"):
        if marker in p:
            p = p[: p.index(marker)].strip().rstrip(".")
            break
    return p[:120] if p else prompt[:80]


# ─────────────────────────────────────────────────────────────────────────────
# PPTX HELPERS
# ─────────────────────────────────────────────────────────────────────────────


def _prs():
    prs = Presentation()
    prs.slide_width = PIn(W)
    prs.slide_height = PIn(H)
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg_white(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PRGB(0xFF, 0xFF, 0xFF)


def _logo_pptx(slide, x=0.12, y=0.08, w=1.25):
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, PIn(x), PIn(y), PIn(w))


def _txt(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=16,
    bold=False,
    color=None,
    font=FONT_BODY,
    align=PP_ALIGN.LEFT,
    italic=False,
):
    txb = slide.shapes.add_textbox(PIn(x), PIn(y), PIn(w), PIn(h))
    txb.text_frame.word_wrap = True
    p = txb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = PPt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color or P_NEGRO


def _rect_shape(slide, x, y, w, h, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(1, PIn(x), PIn(y), PIn(w), PIn(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = PPt(1.5)
    else:
        shape.line.fill.background()
    return shape


def _verde_strip(slide):
    _rect_shape(slide, 0.0, 5.28, W, 0.35, P_VERDE)


def _left_brand_strip(slide):
    """Thin left-side brand strip for visual identity on content slides."""
    _rect_shape(slide, 0.0, 0.0, 0.08, H, P_VERDE)


def _email_tag(slide, email="Misterruddy@laspalmas.edu.bo"):
    _txt(slide, email, 0.18, 0.04, 3.5, 0.28, size=9, italic=True, color=P_GRIS)


def _bullets_block(slide, items, x, y, w, h, size=15):
    if not items:
        return
    txb = slide.shapes.add_textbox(PIn(x), PIn(y), PIn(w), PIn(h))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        s = str(item).strip()
        if not s:
            continue
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_before = PPt(4)
        run = para.add_run()
        run.text = "• " + s.replace("**", "")
        run.font.size = PPt(size)
        run.font.name = FONT_BODY
        run.font.color.rgb = P_NEGRO


def _img_placeholder(slide, x=0.18, y=1.12, w=4.50, h=3.90, prompt=""):
    _rect_shape(slide, x, y, w, h, PRGB(0xF1, 0xF8, 0xE9), P_VERDE)
    _txt(
        slide,
        "🖼 Imagen sugerida",
        x + 0.1,
        y + 1.5,
        w - 0.2,
        0.5,
        size=10,
        italic=True,
        color=PRGB(0x55, 0x88, 0x55),
        align=PP_ALIGN.CENTER,
    )
    if prompt:
        scene = _extract_scene(prompt)
        if scene:
            _txt(
                slide,
                scene,
                x + 0.1,
                y + 2.1,
                w - 0.2,
                1.5,
                size=8,
                italic=True,
                color=PRGB(0x77, 0x99, 0x77),
                align=PP_ALIGN.CENTER,
            )


def _yellow_callout(slide, text, x=0.18, y=3.85, w=4.50, h=1.12):
    _rect_shape(slide, x, y, w, h, P_AMARILLO)
    txb = slide.shapes.add_textbox(
        PIn(x + 0.12), PIn(y + 0.08), PIn(w - 0.24), PIn(h - 0.16)
    )
    txb.text_frame.word_wrap = True
    p = txb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(text).replace("**", "")
    run.font.bold = True
    run.font.size = PPt(13)
    run.font.name = FONT_BODY
    run.font.color.rgb = P_NEGRO


# ── SLIDE LAYOUTS ─────────────────────────────────────────────────────────────


def _slide_cover(prs, unidad, tema):
    sl = _blank(prs)
    _bg_white(sl)
    _verde_strip(sl)
    _rect_shape(sl, 0.0, 0.0, 0.10, H, P_VERDE)
    _logo_pptx(sl, x=0.22, y=0.20, w=2.0)
    _txt(
        sl,
        unidad,
        0.25,
        1.40,
        9.5,
        0.55,
        size=14,
        italic=True,
        color=PRGB(0x2E, 0x7D, 0x32),
    )
    _txt(
        sl,
        tema,
        0.25,
        1.95,
        9.5,
        2.50,
        size=48,
        bold=True,
        font=FONT_TITLE,
        color=P_NEGRO,
    )
    _txt(sl, _fecha_es(), 0.25, 4.35, 7.0, 0.50, size=13, italic=True, color=P_GRIS)


def _slide_content(prs, titulo, concepto, bullets, callout="", prompt=""):
    sl = _blank(prs)
    _bg_white(sl)
    _verde_strip(sl)
    _left_brand_strip(sl)
    _logo_pptx(sl)
    _email_tag(sl)
    titulo_size = 22 if len(titulo) < 30 else 18 if len(titulo) < 45 else 15
    _txt(
        sl,
        titulo,
        0.18,
        0.38,
        9.5,
        0.70,
        size=titulo_size,
        bold=True,
        font=FONT_TITLE,
        color=P_NEGRO,
        align=PP_ALIGN.CENTER,
    )
    _img_placeholder(sl, prompt=prompt)
    if callout:
        _yellow_callout(sl, callout[:120])
    if concepto:
        _txt(sl, concepto, 4.90, 1.12, 5.0, 1.10, size=14, color=P_NEGRO)
    _bullets_block(sl, bullets[:5], 4.90, 2.28, 5.0, 2.80)


def _slide_kinetic(prs, instruccion=""):
    sl = _blank(prs)
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = P_AMARILLO
    _verde_strip(sl)
    _logo_pptx(sl)
    _txt(
        sl,
        "🏃 ¡Pausa Kinestésica!",
        0.5,
        1.5,
        9.0,
        1.0,
        size=36,
        bold=True,
        font=FONT_TITLE,
        color=P_VERDE,
        align=PP_ALIGN.CENTER,
    )
    if instruccion:
        clean = instruccion.replace("**", "")[:150]
        _txt(
            sl,
            clean,
            1.0,
            2.85,
            8.0,
            1.5,
            size=18,
            color=P_NEGRO,
            align=PP_ALIGN.CENTER,
        )


def generar_pptx_diapositivas(res_m1b, tema: str, unidad: str) -> bytes:
    prs = _prs()
    if isinstance(res_m1b, list):
        slides_data = res_m1b
    elif isinstance(res_m1b, dict):
        slides_data = (
            res_m1b.get("diapositivas")
            or res_m1b.get("slides")
            or (list(res_m1b.values())[0] if res_m1b else [])
        )
    else:
        slides_data = []

    _slide_cover(prs, unidad, tema)

    for i, sd in enumerate(slides_data):
        if not isinstance(sd, dict):
            continue
        titulo = str(sd.get("titulo") or sd.get("title") or f"Slide {i + 1}")
        guion = str(sd.get("guion_docente") or sd.get("guion") or "")
        texto = str(sd.get("texto_pantalla") or sd.get("texto") or "")
        prompt = str(sd.get("prompt_imagen") or sd.get("prompt") or "")
        callout = str(sd.get("callout") or "")
        tl = titulo.lower()
        bullets = [
            b.strip().replace("**", "").lstrip("•- ")
            for b in texto.replace("•", "\n").split("\n")
            if b.strip()
        ]

        if any(x in tl for x in ["kinest", "pausa", "break", "movimiento"]):
            _slide_kinetic(prs, guion or texto[:150])
        else:
            concepto = bullets[0] if bullets else ""
            rest = bullets[1:5] if len(bullets) > 1 else []
            _slide_content(
                prs,
                titulo,
                concepto,
                rest,
                callout=callout or guion[:100],
                prompt=prompt,
            )

    # Closing
    sl_fin = _blank(prs)
    _bg_white(sl_fin)
    _verde_strip(sl_fin)
    _logo_pptx(sl_fin, x=3.80, y=1.10, w=2.40)
    _txt(
        sl_fin,
        "¡Excelente trabajo hoy!",
        0.5,
        3.50,
        9.0,
        0.80,
        size=32,
        bold=True,
        font=FONT_TITLE,
        color=P_VERDE,
        align=PP_ALIGN.CENTER,
    )
    _txt(
        sl_fin,
        "Misterruddy@laspalmas.edu.bo",
        0.5,
        4.60,
        9.0,
        0.40,
        size=11,
        italic=True,
        color=P_GRIS,
        align=PP_ALIGN.CENTER,
    )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
