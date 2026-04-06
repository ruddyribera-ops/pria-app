"""
slide_generator.py — Motor Premium de Diapositivas PRIA v5.6
Claude API enriquece contenido → python-pptx genera slides profesionales.
30–50 palabras en pantalla. Referencias a páginas del libro incluidas.
Sin dependencias de Node.js ni pptxgenjs.
"""

import io, json, os, re
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches as PIn, Pt as PPt
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.text import PP_ALIGN

# ─── PATHS & CONSTANTS ───────────────────────────────────────────────────────
_DIR      = Path(__file__).parent
LOGO_PATH = str(_DIR / "logo_laspalmas.png")
if not os.path.exists(LOGO_PATH):
    LOGO_PATH = str(Path.cwd() / "logo_laspalmas.png")

W, H = 10.0, 5.625
FONT_TITLE = "Impact"
FONT_BODY  = "Calibri"

P_VERDE    = PRGB(0x1B, 0x5E, 0x20)
P_VERDE2   = PRGB(0x2E, 0x7D, 0x32)
P_NEGRO    = PRGB(0x1A, 0x1A, 0x1A)
P_BLANCO   = PRGB(0xFF, 0xFF, 0xFF)
P_GRIS     = PRGB(0x55, 0x55, 0x55)
P_AMARILLO = PRGB(0xFF, 0xEE, 0x00)
P_AZUL     = PRGB(0x1A, 0x3A, 0x6A)
P_TEAL     = PRGB(0x00, 0x7A, 0x6E)
P_VERDE_BG = PRGB(0xF1, 0xF8, 0xE9)

_PAGE_RE = re.compile(r'p[áa]g(?:ina)?\.?\s*(\d+)', re.IGNORECASE)

_MESES_ES = {1:"enero",2:"febrero",3:"marzo",4:"abril",5:"mayo",6:"junio",
             7:"julio",8:"agosto",9:"septiembre",10:"octubre",11:"noviembre",12:"diciembre"}

def _fecha_es() -> str:
    d = datetime.now()
    return f"{d.day:02d} de {_MESES_ES[d.month]} de {d.year}"

def _extract_scene(prompt: str) -> str:
    """Extract only the scene (PARTE 2) from the three-part M1b image prompt."""
    p = str(prompt).strip()
    for marker in ("Square format.", "square format."):
        if marker in p:
            p = p[p.index(marker) + len(marker):].strip().lstrip(".")
            break
    for marker in ("Consistent illustration", "consistent illustration"):
        if marker in p:
            p = p[:p.index(marker)].strip().rstrip(".")
            break
    return p[:120] if p else prompt[:80]


# ─── AVAILABILITY ─────────────────────────────────────────────────────────────
def node_disponible() -> bool:
    """Premium engine is pure Python — always available when Anthropic key is set."""
    return True


def instalar_pptxgenjs() -> str:
    return "✅ Motor premium activo (Python nativo + Claude API). No requiere instalación adicional."


# ─── PPTX DRAW HELPERS ────────────────────────────────────────────────────────
def _prs():
    prs = Presentation()
    prs.slide_width  = PIn(W)
    prs.slide_height = PIn(H)
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, r=255, g=255, b=255):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PRGB(r, g, b)


def _rect(slide, x, y, w, h, fill, line=None, lw=1.5):
    shape = slide.shapes.add_shape(1, PIn(x), PIn(y), PIn(w), PIn(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = PPt(lw)
    else:
        shape.line.fill.background()
    return shape


def _txt(slide, text, x, y, w, h, size=14, bold=False, color=None,
         font=FONT_BODY, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(PIn(x), PIn(y), PIn(w), PIn(h))
    txb.text_frame.word_wrap = True
    p   = txb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text            = str(text)
    run.font.size       = PPt(size)
    run.font.bold       = bold
    run.font.italic     = italic
    run.font.name       = font
    run.font.color.rgb  = color or P_NEGRO


def _logo(slide, x=0.12, y=0.06, w=1.20):
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, PIn(x), PIn(y), PIn(w))


def _verde_strip(slide):
    _rect(slide, 0.0, 5.28, W, 0.35, P_VERDE)


def _left_accent(slide, color=None):
    _rect(slide, 0.0, 0.0, 0.10, H, color or P_VERDE)


def _bullets_block(slide, bullets, x, y, w, h, size=13):
    if not bullets:
        return
    txb = slide.shapes.add_textbox(PIn(x), PIn(y), PIn(w), PIn(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        s = str(b).strip().replace("**", "").lstrip("•▸-– ")
        if not s:
            continue
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_before = PPt(5)
        run = para.add_run()
        run.text           = "▸  " + s
        run.font.size      = PPt(size)
        run.font.name      = FONT_BODY
        run.font.color.rgb = P_NEGRO


def _page_badge(slide, page_ref: str):
    """Teal badge top-right corner showing page reference."""
    _rect(slide, 8.58, 0.35, 1.30, 0.32, P_TEAL)
    _txt(slide, page_ref, 8.62, 0.37, 1.22, 0.28,
         size=9, bold=True, color=P_BLANCO, align=PP_ALIGN.CENTER)


def _callout_box(slide, text, x=4.80, y=4.44, w=4.95, h=0.68):
    """Yellow callout / memorable phrase."""
    _rect(slide, x, y, w, h, P_AMARILLO)
    txb = slide.shapes.add_textbox(PIn(x + 0.12), PIn(y + 0.08), PIn(w - 0.24), PIn(h - 0.16))
    txb.text_frame.word_wrap = True
    p   = txb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text           = str(text).replace("**", "")[:200]
    run.font.bold      = True
    run.font.size      = PPt(12)
    run.font.name      = FONT_BODY
    run.font.color.rgb = P_NEGRO


def _img_zone(slide, x=0.18, y=1.12, w=4.45, h=3.48, prompt=""):
    """Image placeholder with optional prompt hint (scene only, no style prefix)."""
    _rect(slide, x, y, w, h, P_VERDE_BG, P_VERDE2, lw=0.75)
    _txt(slide, "🖼  Imagen sugerida", x + 0.1, y + h / 2 - 0.22, w - 0.2, 0.40,
         size=10, italic=True, color=PRGB(0x55, 0x88, 0x55), align=PP_ALIGN.CENTER)
    if prompt:
        scene = _extract_scene(prompt)
        if scene:
            _txt(slide, scene, x + 0.1, y + h / 2 + 0.15, w - 0.2, 0.70,
                 size=7, italic=True, color=PRGB(0x77, 0x99, 0x77), align=PP_ALIGN.CENTER)


# ─── SLIDE LAYOUTS ────────────────────────────────────────────────────────────
def _slide_num_tag(slide, label: str):
    """Small slide counter tag in the green strip bottom-right."""
    _txt(slide, label, 9.00, 5.28, 0.90, 0.34,
         size=8, color=P_BLANCO, align=PP_ALIGN.RIGHT)


def _slide_cover(prs, unidad: str, tema: str, total: int = 0):
    sl = _blank(prs)
    _bg(sl)
    _verde_strip(sl)
    _rect(sl, 0.0, 0.0, 0.12, H, P_VERDE)
    _logo(sl, x=0.24, y=0.18, w=2.0)
    _txt(sl, unidad, 0.28, 1.35, 9.5, 0.50,
         size=13, italic=True, color=P_VERDE2)
    _txt(sl, tema,   0.28, 1.85, 9.5, 2.50,
         size=44, bold=True, font=FONT_TITLE, color=P_NEGRO)
    _txt(sl, _fecha_es(), 0.28, 4.30, 7.0, 0.50,
         size=12, italic=True, color=P_GRIS)
    if total > 0:
        _slide_num_tag(sl, f"1 / {total + 2}")


def _slide_content_premium(prs, titulo: str, concepto: str, bullets: list,
                            callout: str = "", prompt: str = "",
                            page_ref: str = "", slide_num: int = 0, total: int = 0):
    """
    Premium two-zone layout:
      Left  (0.18–4.63): image placeholder
      Right (4.80–9.75): concept sentence + bullets + callout
    Page reference badge top-right. Slide counter in green strip.
    """
    sl = _blank(prs)
    _bg(sl)
    _verde_strip(sl)
    _left_accent(sl)
    _logo(sl)

    # Email
    _txt(sl, "Misterruddy@laspalmas.edu.bo", 0.18, 0.04, 3.5, 0.26,
         size=8, italic=True, color=P_GRIS)

    # Title — auto-size
    ts = 22 if len(titulo) < 30 else 18 if len(titulo) < 50 else 14
    _txt(sl, titulo, 0.18, 0.36, 9.55, 0.72,
         size=ts, bold=True, font=FONT_TITLE, color=P_NEGRO, align=PP_ALIGN.CENTER)

    # Left zone: image
    _img_zone(sl, prompt=prompt)

    # Right zone: concept + bullets
    RX, RW = 4.80, 4.95
    ry = 1.14

    if concepto:
        txb = sl.shapes.add_textbox(PIn(RX), PIn(ry), PIn(RW), PIn(1.30))
        txb.text_frame.word_wrap = True
        p   = txb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text           = str(concepto).replace("**", "")
        run.font.size      = PPt(14)
        run.font.name      = FONT_BODY
        run.font.color.rgb = P_NEGRO
        ry += 1.35

    if bullets:
        available_h = 4.42 - ry
        _bullets_block(sl, bullets[:5], RX, ry, RW, max(available_h, 1.0), size=13)

    # Callout
    if callout:
        _callout_box(sl, callout)

    # Page reference badge
    if page_ref:
        _page_badge(sl, page_ref)

    # Slide counter
    if total > 0:
        _txt(sl, f"{slide_num} / {total}", 9.00, 5.28, 0.90, 0.34,
             size=8, color=P_BLANCO, align=PP_ALIGN.RIGHT)


def _slide_kinetic(prs, instruccion: str = ""):
    sl = _blank(prs)
    _bg(sl, 255, 238, 0)
    _verde_strip(sl)
    _logo(sl)
    _txt(sl, "🏃 ¡Pausa Kinestésica!", 0.5, 1.45, 9.0, 1.0,
         size=36, bold=True, font=FONT_TITLE, color=P_VERDE, align=PP_ALIGN.CENTER)
    if instruccion:
        _txt(sl, instruccion.replace("**", "")[:200], 1.0, 2.75, 8.0, 1.5,
             size=17, color=P_NEGRO, align=PP_ALIGN.CENTER)


def _slide_closing(prs, email: str = "Misterruddy@laspalmas.edu.bo", total: int = 0):
    sl = _blank(prs)
    _bg(sl)
    _verde_strip(sl)
    _logo(sl, x=3.80, y=1.10, w=2.40)
    _txt(sl, "¡Excelente trabajo hoy!", 0.5, 3.45, 9.0, 0.80,
         size=32, bold=True, font=FONT_TITLE, color=P_VERDE, align=PP_ALIGN.CENTER)
    _txt(sl, email, 0.5, 4.55, 9.0, 0.40,
         size=11, italic=True, color=P_GRIS, align=PP_ALIGN.CENTER)
    if total > 0:
        _slide_num_tag(sl, f"{total + 2} / {total + 2}")


# ─── CLAUDE ENRICHMENT ────────────────────────────────────────────────────────
def _extract_page_ref(text: str) -> str:
    m = _PAGE_RE.search(str(text))
    return f"Pág. {m.group(1)}" if m else ""


def _is_kinetic(titulo: str) -> bool:
    kw = ["kinest", "pausa", "break", "movimiento", "activ"]
    return any(k in titulo.lower() for k in kw)


def _enrich_with_claude(client, slides_raw: list, tema: str, unidad: str) -> list:
    """
    Single Claude API call — enriches all content slides in one shot.
    Returns list of dicts ready for rendering.
    """
    content_slides = [
        {
            "idx": i,
            "titulo":  s.get("titulo") or s.get("title") or f"Slide {i+1}",
            "guion":   s.get("guion_docente") or s.get("guion") or "",
            "texto":   s.get("texto_pantalla") or s.get("texto") or "",
            "callout": s.get("callout") or "",
        }
        for i, s in enumerate(slides_raw)
        if isinstance(s, dict) and not _is_kinetic(s.get("titulo") or s.get("title") or "")
    ]

    if not content_slides:
        return []

    prompt_text = f"""Eres un diseñador instruccional experto para 5to de Primaria (10-11 años).

Lección: "{tema}" | Unidad: "{unidad}"

Para cada diapositiva genera:
- **concepto**: 1-2 oraciones (máx 40 palabras), lenguaje claro para niños. Explica el "qué" y el "por qué" de forma directa.
- **bullets**: 3-4 puntos clave, máx 9 palabras cada uno. Concretos y accionables.
- **callout**: 1 frase memorable o pregunta de reflexión, máx 15 palabras.
- **page_ref**: Si el guion menciona páginas ("página 23", "pág. 45", "p.67"), extrae "Pág. XX". Si no hay, pon "".

REGLAS CRÍTICAS:
- concepto + bullets juntos NO deben superar 50 palabras en pantalla
- Lenguaje directo, sin jerga académica
- Los bullets empiezan con verbo en infinitivo o sustantivo concreto

DIAPOSITIVAS:
{json.dumps(content_slides, ensure_ascii=False, indent=2)}

Responde SOLO con JSON. Sin texto extra. Formato exacto:
{{
  "slides": [
    {{
      "idx": 0,
      "concepto": "...",
      "bullets": ["...", "...", "..."],
      "callout": "...",
      "page_ref": ""
    }}
  ]
}}"""

    # Cache the large instruction block to reduce cost ~90% on repeated calls
    response = client.messages.create(
        model="claude-opus-4-6",
        max_tokens=4096,
        messages=[{
            "role": "user",
            "content": [{
                "type": "text",
                "text": prompt_text,
                "cache_control": {"type": "ephemeral"}
            }]
        }]
    )

    raw = response.content[0].text.strip()
    m   = re.search(r'\{[\s\S]*\}', raw)
    if not m:
        raise ValueError(f"Claude no devolvió JSON válido. Respuesta: {raw[:300]}")

    enriched_map = {}
    for item in json.loads(m.group(0)).get("slides", []):
        enriched_map[item["idx"]] = item

    return enriched_map


def _build_slide_list(slides_raw: list, enriched_map: dict) -> list:
    """Merge raw slide data with Claude enrichment."""
    result = []
    for i, sd in enumerate(slides_raw):
        if not isinstance(sd, dict):
            continue
        titulo = sd.get("titulo") or sd.get("title") or f"Slide {i+1}"
        prompt = sd.get("prompt_imagen") or sd.get("prompt") or ""
        guion  = sd.get("guion_docente") or sd.get("guion") or ""

        if _is_kinetic(titulo):
            result.append({"_type": "kinetic", "titulo": titulo, "guion": guion})
            continue

        enr      = enriched_map.get(i, {})
        page_ref = enr.get("page_ref") or _extract_page_ref(guion)

        result.append({
            "_type":   "content",
            "titulo":  titulo,
            "concepto": enr.get("concepto") or sd.get("texto_pantalla") or sd.get("texto") or "",
            "bullets": enr.get("bullets") or [],
            "callout": enr.get("callout") or sd.get("callout") or "",
            "page_ref": page_ref,
            "prompt":  prompt,
            "guion":   guion,
        })
    return result


# ─── PUBLIC API ───────────────────────────────────────────────────────────────
def generar_pptx_con_claude(slides_data, tema: str, unidad: str,
                             anthropic_api_key) -> bytes:
    """
    Premium PPTX generator.

    1. Calls Claude API to rewrite slide text: 30-50 words per screen,
       page references, memorable callouts — all calibrated for age 10-11.
    2. Renders with python-pptx using a professional two-zone layout.

    Args:
        slides_data:       List of slide dicts (or dict wrapping a list).
        tema:              Active topic string.
        unidad:            Unit title string.
        anthropic_api_key: API key string or list (first element used).

    Returns:
        bytes: PPTX file content.
    """
    import anthropic

    # Handle key as list or string
    if isinstance(anthropic_api_key, (list, tuple)):
        anthropic_api_key = anthropic_api_key[0]
    api_key = str(anthropic_api_key).strip()

    # Normalize slides_data to a list
    if isinstance(slides_data, dict):
        slides_data = (slides_data.get("diapositivas")
                       or slides_data.get("slides")
                       or (list(slides_data.values())[0] if slides_data else []))
    if not isinstance(slides_data, list):
        slides_data = []

    client = anthropic.Anthropic(api_key=api_key)

    # Enrich with Claude
    enriched_map = _enrich_with_claude(client, slides_data, tema, unidad)

    # Merge raw + enriched
    slide_list = _build_slide_list(slides_data, enriched_map)

    # Count content slides for numbering
    total_content = sum(1 for s in slide_list if s["_type"] == "content")

    # Render
    prs     = _prs()
    _slide_cover(prs, unidad, tema, total=total_content)

    content_num = 0
    for s in slide_list:
        if s["_type"] == "kinetic":
            _slide_kinetic(prs, s.get("guion", "")[:180])
        else:
            content_num += 1
            _slide_content_premium(
                prs,
                titulo    = s["titulo"],
                concepto  = s["concepto"],
                bullets   = s["bullets"],
                callout   = s["callout"],
                prompt    = s["prompt"],
                page_ref  = s["page_ref"],
                slide_num = content_num,
                total     = total_content,
            )

    _slide_closing(prs, total=total_content)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
