"""
exportar.py — Módulo de exportación PRIA v5.6
Plantilla oficial Las Palmas School — Listo para imprimir.

Fixes v5.6:
  - Asteriscos dobles eliminados (inline bold parseado correctamente)
  - Guiones "---" ya no aparecen como texto
  - IA ya no se autopresenta en los documentos
  - Header "Fichas de Trabajo" sin recorte
  - Ficha Gamificada rediseñada como ficha de trabajo real
  - Espaciado generoso — documentos respiran
  - Sopa de letras real generada con cuadrícula
  - Espacios de respuesta para el alumno
"""

import io, json, os, re, random
from datetime import datetime
from pathlib import Path

from docx import Document
from docx.shared import Pt, RGBColor, Inches, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_ALIGN_VERTICAL
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches as PIn, Pt as PPt
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.text import PP_ALIGN

# ─── LOGO PATH ────────────────────────────────────────────────────────────────
_DIR      = Path(__file__).parent
LOGO_PATH = str(_DIR / "logo_laspalmas.png")
if not os.path.exists(LOGO_PATH):
    LOGO_PATH = str(Path.cwd() / "logo_laspalmas.png")

def _logo_stream():
    with open(LOGO_PATH, "rb") as f:
        return io.BytesIO(f.read())

# ─── COLORES ─────────────────────────────────────────────────────────────────
P_VERDE    = PRGB(0x1B, 0x5E, 0x20)
P_NEGRO    = PRGB(0x1A, 0x1A, 0x1A)
P_BLANCO   = PRGB(0xFF, 0xFF, 0xFF)
P_GRIS     = PRGB(0x55, 0x55, 0x55)
P_AMARILLO = PRGB(0xFF, 0xEE, 0x00)

D_VERDE    = RGBColor(0x1B, 0x5E, 0x20)
D_VERDE2   = RGBColor(0x2E, 0x7D, 0x32)
D_NEGRO    = RGBColor(0x1A, 0x1A, 0x1A)
D_BLANCO   = RGBColor(0xFF, 0xFF, 0xFF)
D_ROJO     = RGBColor(0xCC, 0x22, 0x00)
D_AZUL     = RGBColor(0x1A, 0x3A, 0x6A)
D_TEAL     = RGBColor(0x00, 0x7A, 0x6E)
D_MORADO   = RGBColor(0x6A, 0x3D, 0x9A)
D_NARANJA  = RGBColor(0xCC, 0x55, 0x00)
D_GRIS     = RGBColor(0x44, 0x44, 0x44)
D_GRIS_CLR = RGBColor(0xAA, 0xAA, 0xAA)

FONT_TITLE = "Impact"
FONT_BODY  = "Calibri"

_MESES_ES = {1:"enero",2:"febrero",3:"marzo",4:"abril",5:"mayo",6:"junio",
             7:"julio",8:"agosto",9:"septiembre",10:"octubre",11:"noviembre",12:"diciembre"}

def _fecha_es() -> str:
    d = datetime.now()
    return f"{d.day:02d} de {_MESES_ES[d.month]} de {d.year}"

def _extract_scene(prompt: str) -> str:
    """Extract only PARTE 2 (the scene) from the three-part M1b image prompt.
    Strips the style prefix (PARTE 1) and technical suffix (PARTE 3)."""
    p = str(prompt).strip()
    # Remove PARTE 1 — ends with "Square format."
    for marker in ("Square format.", "square format."):
        if marker in p:
            p = p[p.index(marker) + len(marker):].strip().lstrip(".")
            break
    # Remove PARTE 3 — starts with "Consistent illustration"
    for marker in ("Consistent illustration", "consistent illustration"):
        if marker in p:
            p = p[:p.index(marker)].strip().rstrip(".")
            break
    return p[:120] if p else prompt[:80]

# ─── FILTRO: líneas de autopresentación de la IA ─────────────────────────────
_AI_INTRO_FRAGMENTS = [
    "soy el motor", "motor m0", "motor m1", "motor m2",
    "método palma-ribera", "palma-ribera", "he diseñado este",
    "me encanta acompañarte", "como experto en currículo",
    "bienvenidos, futuros",
    "a continuación te presento", "como motor", "he sido diseñado",
    "este pop quiz", "esta evaluación ha sido", "bienvenido a la",
    "estimado docente,", "querido estudiante,", "hola, soy",
    "aquí tienes", "a continuación encontrarás",
]

def _es_autopresentacion(linea: str) -> bool:
    ll = linea.lower()[:120]
    return any(f in ll for f in _AI_INTRO_FRAGMENTS)


# ─────────────────────────────────────────────────────────────────────────────
# PPTX HELPERS (fallback python-pptx)
# ─────────────────────────────────────────────────────────────────────────────
W, H = 10.0, 5.625

def _prs():
    prs = Presentation()
    prs.slide_width  = PIn(W)
    prs.slide_height = PIn(H)
    return prs

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _bg_white(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PRGB(0xFF,0xFF,0xFF)

def _logo_pptx(slide, x=0.12, y=0.08, w=1.25):
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, PIn(x), PIn(y), PIn(w))

def _txt(slide, text, x, y, w, h, size=16, bold=False, color=None,
         font=FONT_BODY, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(PIn(x), PIn(y), PIn(w), PIn(h))
    txb.text_frame.word_wrap = True
    p   = txb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text); run.font.size = PPt(size); run.font.bold = bold
    run.font.italic = italic; run.font.name = font
    run.font.color.rgb = color or P_NEGRO

def _rect_shape(slide, x, y, w, h, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(1, PIn(x), PIn(y), PIn(w), PIn(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb; shape.line.width = PPt(1.5)
    else:
        shape.line.fill.background()
    return shape

def _verde_strip(slide):
    _rect_shape(slide, 0.0, 5.28, W, 0.35, P_VERDE)

def _left_brand_strip(slide):
    """Thin left-side brand strip for visual identity on content slides."""
    _rect_shape(slide, 0.0, 0.0, 0.08, H, P_VERDE)

def _email_tag(slide, email="Misterruddy@laspalmas.edu.bo"):
    _txt(slide, email, 0.18, 0.04, 3.5, 0.28, size=9, italic=True, color=P_GRIS)

def _bullets_block(slide, items, x, y, w, h, size=15):
    if not items: return
    txb = slide.shapes.add_textbox(PIn(x), PIn(y), PIn(w), PIn(h))
    tf  = txb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        s = str(item).strip()
        if not s: continue
        para = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        para.space_before = PPt(4)
        run = para.add_run()
        run.text = "• " + s.replace("**","")
        run.font.size = PPt(size); run.font.name = FONT_BODY
        run.font.color.rgb = P_NEGRO

def _img_placeholder(slide, x=0.18, y=1.12, w=4.50, h=3.90, prompt=""):
    _rect_shape(slide, x, y, w, h, PRGB(0xF1,0xF8,0xE9), P_VERDE)
    _txt(slide, "🖼 Imagen sugerida", x+0.1, y+1.5, w-0.2, 0.5,
         size=10, italic=True, color=PRGB(0x55,0x88,0x55), align=PP_ALIGN.CENTER)
    if prompt:
        scene = _extract_scene(prompt)
        if scene:
            _txt(slide, scene, x+0.1, y+2.1, w-0.2, 1.5,
                 size=8, italic=True, color=PRGB(0x77,0x99,0x77), align=PP_ALIGN.CENTER)

def _yellow_callout(slide, text, x=0.18, y=3.85, w=4.50, h=1.12):
    _rect_shape(slide, x, y, w, h, P_AMARILLO)
    txb = slide.shapes.add_textbox(PIn(x+0.12), PIn(y+0.08), PIn(w-0.24), PIn(h-0.16))
    txb.text_frame.word_wrap = True
    p = txb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(text).replace("**","")
    run.font.bold = True; run.font.size = PPt(13)
    run.font.name = FONT_BODY; run.font.color.rgb = P_NEGRO

# ── SLIDE LAYOUTS ─────────────────────────────────────────────────────────────
def _slide_cover(prs, unidad, tema):
    sl = _blank(prs); _bg_white(sl); _verde_strip(sl)
    _rect_shape(sl, 0.0, 0.0, 0.10, H, P_VERDE)
    _logo_pptx(sl, x=0.22, y=0.20, w=2.0)
    _txt(sl, unidad, 0.25, 1.40, 9.5, 0.55,
         size=14, italic=True, color=PRGB(0x2E,0x7D,0x32))
    _txt(sl, tema,   0.25, 1.95, 9.5, 2.50,
         size=48, bold=True, font=FONT_TITLE, color=P_NEGRO)
    _txt(sl, _fecha_es(), 0.25, 4.35, 7.0, 0.50,
         size=13, italic=True, color=P_GRIS)

def _slide_content(prs, titulo, concepto, bullets, callout="", prompt=""):
    sl = _blank(prs); _bg_white(sl); _verde_strip(sl)
    _left_brand_strip(sl)
    _logo_pptx(sl); _email_tag(sl)
    # Title — fixed height, smaller font to prevent overflow
    titulo_size = 22 if len(titulo) < 30 else 18 if len(titulo) < 45 else 15
    _txt(sl, titulo, 0.18, 0.38, 9.5, 0.70,
         size=titulo_size, bold=True, font=FONT_TITLE,
         color=P_NEGRO, align=PP_ALIGN.CENTER)
    # Image left
    _img_placeholder(sl, prompt=prompt)
    if callout:
        _yellow_callout(sl, callout[:120])
    # Content right — starts well below logo, above green strip
    if concepto:
        _txt(sl, concepto, 4.90, 1.12, 5.0, 1.10, size=14, color=P_NEGRO)
    _bullets_block(sl, bullets[:5], 4.90, 2.28, 5.0, 2.80)

def _slide_kinetic(prs, instruccion=""):
    sl = _blank(prs)
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = P_AMARILLO
    _verde_strip(sl); _logo_pptx(sl)
    _txt(sl, "🏃 ¡Pausa Kinestésica!", 0.5, 1.5, 9.0, 1.0,
         size=36, bold=True, font=FONT_TITLE, color=P_VERDE, align=PP_ALIGN.CENTER)
    if instruccion:
        clean = instruccion.replace("**","")[:150]
        _txt(sl, clean, 1.0, 2.85, 8.0, 1.5,
             size=18, color=P_NEGRO, align=PP_ALIGN.CENTER)

def generar_pptx_diapositivas(res_m1b, tema: str, unidad: str) -> bytes:
    prs = _prs()
    if isinstance(res_m1b, list):
        slides_data = res_m1b
    elif isinstance(res_m1b, dict):
        slides_data = (res_m1b.get("diapositivas") or res_m1b.get("slides")
                       or (list(res_m1b.values())[0] if res_m1b else []))
    else:
        slides_data = []

    _slide_cover(prs, unidad, tema)

    for i, sd in enumerate(slides_data):
        if not isinstance(sd, dict): continue
        titulo  = str(sd.get("titulo") or sd.get("title") or f"Slide {i+1}")
        guion   = str(sd.get("guion_docente") or sd.get("guion") or "")
        texto   = str(sd.get("texto_pantalla") or sd.get("texto") or "")
        prompt  = str(sd.get("prompt_imagen") or sd.get("prompt") or "")
        callout = str(sd.get("callout") or "")
        tl = titulo.lower()
        bullets = [b.strip().replace("**","").lstrip("•- ") 
                   for b in texto.replace("•","\n").split("\n") if b.strip()]

        if any(x in tl for x in ["kinest","pausa","break","movimiento"]):
            _slide_kinetic(prs, guion or texto[:150])
        else:
            concepto = bullets[0] if bullets else ""
            rest     = bullets[1:5] if len(bullets)>1 else []
            _slide_content(prs, titulo, concepto, rest,
                           callout=callout or guion[:100], prompt=prompt)

    # Closing
    sl_fin = _blank(prs); _bg_white(sl_fin); _verde_strip(sl_fin)
    _logo_pptx(sl_fin, x=3.80, y=1.10, w=2.40)
    _txt(sl_fin, "¡Excelente trabajo hoy!", 0.5, 3.50, 9.0, 0.80,
         size=32, bold=True, font=FONT_TITLE, color=P_VERDE, align=PP_ALIGN.CENTER)
    _txt(sl_fin, "Misterruddy@laspalmas.edu.bo", 0.5, 4.60, 9.0, 0.40,
         size=11, italic=True, color=P_GRIS, align=PP_ALIGN.CENTER)

    buf = io.BytesIO(); prs.save(buf); return buf.getvalue()


# ─────────────────────────────────────────────────────────────────────────────
# DOCX HELPERS
# ─────────────────────────────────────────────────────────────────────────────
def _nuevo_doc() -> Document:
    doc = Document()
    for s in doc.sections:
        s.top_margin    = Inches(0.75)
        s.bottom_margin = Inches(0.75)
        s.left_margin   = Inches(0.90)
        s.right_margin  = Inches(0.90)
    doc.styles["Normal"].font.name = FONT_BODY
    doc.styles["Normal"].font.size = Pt(11)
    return doc

def _cell_bg(cell, hex6: str):
    tc = cell._tc; tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear"); shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex6); tcPr.append(shd)

def _row_h(row, cm: float):
    tr = row._tr; trPr = tr.get_or_add_trPr()
    h = OxmlElement("w:trHeight")
    h.set(qn("w:val"), str(int(cm * 567))); h.set(qn("w:hRule"), "exact")
    trPr.append(h)

def _runs_from_inline_md(para, text: str, size=11, color=None, italic=False):
    """Add runs to paragraph, converting **bold** inline markdown to actual bold."""
    col = color or D_NEGRO
    parts = re.split(r'\*\*(.+?)\*\*', str(text))
    for i, part in enumerate(parts):
        if not part: continue
        run = para.add_run(part)
        run.font.name      = FONT_BODY
        run.font.size      = Pt(size)
        run.font.color.rgb = col
        run.font.italic    = italic
        if i % 2 == 1:   # odd = the bold part
            run.font.bold = True

def _worksheet_header(doc: Document, materia: str, tema: str, grado: str = "5to"):
    """Navy banner header — logo centred, no text clipping."""
    tbl = doc.add_table(rows=2, cols=3)
    tbl.style = "Table Grid"
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER

    r0 = tbl.rows[0]
    _row_h(r0, 1.5)   # FIX: was 1.0 — text clipped at 1.0

    c0 = r0.cells[0]; _cell_bg(c0, "1A3A6A"); c0.width = Inches(3.2)
    p0 = c0.paragraphs[0]; p0.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p0.paragraph_format.space_before = Pt(6)
    run = p0.add_run("Fichas de Trabajo\n")
    run.font.name = FONT_BODY; run.font.bold = True
    run.font.size = Pt(14); run.font.color.rgb = D_BLANCO
    run2 = p0.add_run("Las Palmas School")
    run2.font.name = FONT_BODY; run2.font.size = Pt(10)
    run2.font.color.rgb = RGBColor(0xC8,0xE6,0xC9)

    c1 = r0.cells[1]; _cell_bg(c1, "1A3A6A"); c1.width = Inches(2.6)
    p1 = c1.paragraphs[0]; p1.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if os.path.exists(LOGO_PATH):
        p1.add_run().add_picture(LOGO_PATH, width=Inches(1.55))

    c2 = r0.cells[2]; _cell_bg(c2, "1A3A6A"); c2.width = Inches(1.7)
    p2 = c2.paragraphs[0]; p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p2.paragraph_format.space_before = Pt(8)
    rg = p2.add_run(f"{grado}\nPRIMARIA")
    rg.font.name = FONT_BODY; rg.font.bold = True
    rg.font.size = Pt(16); rg.font.color.rgb = D_BLANCO

    r1 = tbl.rows[1]; _row_h(r1, 0.6)
    r1.cells[0].merge(r1.cells[2])
    c_sub = r1.cells[0]; _cell_bg(c_sub, "143060")
    ps = c_sub.paragraphs[0]; ps.alignment = WD_ALIGN_PARAGRAPH.CENTER
    rs = ps.add_run(materia.upper())
    rs.font.name = FONT_BODY; rs.font.bold = True
    rs.font.size = Pt(12); rs.font.color.rgb = D_BLANCO

    doc.add_paragraph()

    # Red bold title with underline
    p_t = doc.add_paragraph()
    p_t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_t.paragraph_format.space_before = Pt(4)
    p_t.paragraph_format.space_after  = Pt(6)
    rt = p_t.add_run(tema.upper())
    rt.font.name = FONT_TITLE; rt.font.bold = True
    rt.font.size = Pt(20); rt.font.color.rgb = D_ROJO
    pPr = p_t._p.get_or_add_pPr()
    pBdr = OxmlElement("w:pBdr"); bot = OxmlElement("w:bottom")
    bot.set(qn("w:val"), "single"); bot.set(qn("w:sz"), "12")
    bot.set(qn("w:space"), "4"); bot.set(qn("w:color"), "CC2200")
    pBdr.append(bot); pPr.append(pBdr)
    doc.add_paragraph()

def _section_banner(doc: Document, texto: str, hex_color: str, size=12):
    tbl = doc.add_table(rows=1, cols=1)
    tbl.style = "Table Grid"; tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    cell = tbl.rows[0].cells[0]; _cell_bg(cell, hex_color); _row_h(tbl.rows[0], 0.55)
    p = cell.paragraphs[0]; p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(texto)
    run.font.name = FONT_BODY; run.font.bold = True
    run.font.size = Pt(size); run.font.color.rgb = D_BLANCO
    doc.add_paragraph()

def _h(doc, text, lvl=2, color=None):
    clean = str(text).replace("**","").replace("__","")
    p = doc.add_heading(clean, level=lvl)
    for run in p.runs:
        run.font.color.rgb = color or D_AZUL
        run.font.name = FONT_BODY
        run.font.size = Pt({1:18,2:14,3:12}.get(lvl, 11))
    p.paragraph_format.space_before = Pt(10)
    p.paragraph_format.space_after  = Pt(5)

def _p(doc, text, bold=False, italic=False, color=None, size=11):
    clean = str(text).replace("**","").replace("__","")
    p = doc.add_paragraph()
    run = p.add_run(clean)
    run.font.name = FONT_BODY; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color or D_NEGRO
    p.paragraph_format.space_after  = Pt(5)
    p.paragraph_format.space_before = Pt(3)

def _p_md(doc, text, size=11, color=None):
    """Paragraph with inline **bold** markdown support."""
    p = doc.add_paragraph()
    p.paragraph_format.space_after  = Pt(5)
    p.paragraph_format.space_before = Pt(3)
    _runs_from_inline_md(p, text, size=size, color=color)

def _bul(doc, text, color=None, size=11):
    """Bullet item with inline **bold** support."""
    p = doc.add_paragraph(style="List Bullet")
    p.paragraph_format.space_after  = Pt(4)
    p.paragraph_format.left_indent  = Inches(0.25)
    _runs_from_inline_md(p, str(text), size=size, color=color)

def _bloom_badge_table(doc: Document, verbos: list):
    """Horizontal colored badge row — one cell per Bloom verb."""
    if not verbos: return
    colors = ["1B5E20","007A6E","1A3A6A","6A3D9A","CC5500","2E6DA4","8B0000"]
    tbl = doc.add_table(rows=1, cols=len(verbos))
    tbl.style = "Table Grid"
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    _row_h(tbl.rows[0], 0.52)
    for i, verb in enumerate(verbos):
        cell = tbl.rows[0].cells[i]
        _cell_bg(cell, colors[i % len(colors)])
        p = cell.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_before = Pt(3)
        run = p.add_run(str(verb).upper().replace("**","")[:20])
        run.font.bold  = True
        run.font.name  = FONT_BODY
        run.font.size  = Pt(10)
        run.font.color.rgb = D_BLANCO
    doc.add_paragraph()

def _secuencia_timeline(doc: Document, bloques: list):
    """Visual timeline: coloured left block = phase/time, right = content."""
    block_colors = {"inicio":"1B5E20","desarrollo":"1A3A6A","cierre":"CC5500"}
    for bloque in bloques:
        nombre   = str(bloque.get("nombre","Bloque"))
        duracion = str(bloque.get("duracion","?"))
        objetivo = str(bloque.get("objetivo",""))
        nota     = str(bloque.get("nota",""))

        color = "1A3A6A"
        for key, c in block_colors.items():
            if key in nombre.lower(): color = c; break

        tbl = doc.add_table(rows=1, cols=2)
        tbl.style = "Table Grid"
        tbl.alignment = WD_TABLE_ALIGNMENT.CENTER

        left  = tbl.rows[0].cells[0]
        right = tbl.rows[0].cells[1]
        left.width  = Inches(1.1)
        right.width = Inches(5.9)

        _cell_bg(left, color)

        pl = left.paragraphs[0]
        pl.alignment = WD_ALIGN_PARAGRAPH.CENTER
        pl.paragraph_format.space_before = Pt(4)
        rl1 = pl.add_run(nombre.upper()[:10])
        rl1.font.bold = True; rl1.font.size = Pt(9)
        rl1.font.color.rgb = D_BLANCO; rl1.font.name = FONT_BODY
        rl2 = pl.add_run(f"\n{duracion} min")
        rl2.font.size = Pt(8); rl2.font.name = FONT_BODY
        rl2.font.color.rgb = RGBColor(0xC8,0xE6,0xC9)

        _cell_bg(right, "FAFAFA")
        pr = right.paragraphs[0]
        pr.paragraph_format.left_indent  = Inches(0.10)
        pr.paragraph_format.space_before = Pt(4)
        pr.paragraph_format.space_after  = Pt(3)
        _runs_from_inline_md(pr, objetivo, size=11)

        if nota:
            pn = right.add_paragraph()
            pn.paragraph_format.left_indent = Inches(0.10)
            pn.paragraph_format.space_after = Pt(4)
            rn = pn.add_run(f"💡 {nota.replace('**','')}")
            rn.font.italic = True; rn.font.size = Pt(10)
            rn.font.color.rgb = D_TEAL; rn.font.name = FONT_BODY

        doc.add_paragraph()

def _dua_framework_header(doc: Document):
    """3-column DUA principle header banner."""
    tbl = doc.add_table(rows=1, cols=3)
    tbl.style = "Table Grid"
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    _row_h(tbl.rows[0], 0.55)
    principles = [
        ("🧩 Representación",    "6A3D9A"),
        ("✍️ Acción y Expresión", "1A3A6A"),
        ("💡 Motivación",        "007A6E"),
    ]
    for i, (title, color) in enumerate(principles):
        cell = tbl.rows[0].cells[i]
        _cell_bg(cell, color)
        p = cell.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_before = Pt(4)
        run = p.add_run(title)
        run.font.bold = True; run.font.size = Pt(10)
        run.font.color.rgb = D_BLANCO; run.font.name = FONT_BODY
    doc.add_paragraph()

def _dua_note(doc: Document, icon: str, text: str, icon_hex: str = "007A6E"):
    """Margin-style DUA callout: colored icon badge on left, content on right."""
    tbl = doc.add_table(rows=1, cols=2)
    tbl.style = "Table Grid"
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    left  = tbl.rows[0].cells[0]
    right = tbl.rows[0].cells[1]
    left.width = Inches(0.50)
    _cell_bg(left, icon_hex)
    _row_h(tbl.rows[0], 0.65)
    pl = left.paragraphs[0]
    pl.alignment = WD_ALIGN_PARAGRAPH.CENTER
    pl.paragraph_format.space_before = Pt(3)
    rl = pl.add_run(icon)
    rl.font.size = Pt(11); rl.font.color.rgb = D_BLANCO; rl.font.bold = True
    _cell_bg(right, "F9FAFB")
    pr = right.paragraphs[0]
    pr.paragraph_format.left_indent  = Inches(0.08)
    pr.paragraph_format.space_before = Pt(3)
    pr.paragraph_format.space_after  = Pt(3)
    _runs_from_inline_md(pr, str(text).lstrip("- •"), size=10)
    doc.add_paragraph()

def _tbl(doc, filas: list, hdr_color="1A3A6A"):
    if not filas: return
    cols = list(filas[0].keys())
    tbl = doc.add_table(rows=1, cols=len(cols))
    tbl.style = "Table Grid"
    hdr = tbl.rows[0].cells
    for i, col in enumerate(cols):
        hdr[i].text = col; _cell_bg(hdr[i], hdr_color)
        r = hdr[i].paragraphs[0].runs[0]
        r.font.color.rgb = D_BLANCO; r.font.bold = True
        r.font.name = FONT_BODY; r.font.size = Pt(10)
    for fila in filas:
        row = tbl.add_row().cells
        for i, col in enumerate(cols):
            row[i].text = str(fila.get(col, ""))
            _cell_bg(row[i], "F1F8E9")
            r = row[i].paragraphs[0].runs[0]
            r.font.color.rgb = D_NEGRO; r.font.name = FONT_BODY; r.font.size = Pt(10)
    doc.add_paragraph()

def _md(doc, texto: str):
    """Convert markdown text to DOCX, with all bugs fixed."""
    if not texto: return
    buf = []
    def flush():
        if not buf: return
        validas = [l for l in buf if not all(c in "-| " for c in l)]
        if not validas: buf.clear(); return
        cols = [c.strip() for c in validas[0].split("|") if c.strip()]
        rows = []
        for rl in validas[1:]:
            vals = [c.strip() for c in rl.split("|") if c.strip()]
            if vals: rows.append(dict(zip(cols, vals+[""]*(max(0,len(cols)-len(vals))))))
        if rows: _tbl(doc, rows)
        buf.clear()

    for line in texto.split("\n"):
        s = line.strip()
        if s.startswith("|"): buf.append(s); continue
        else: flush()

        # Skip AI self-introduction lines
        if _es_autopresentacion(s): continue
        # Skip markdown horizontal rules
        if s in ("---", "___", "***") or (len(s)>=3 and all(c=='-' for c in s)): continue
        if not s: doc.add_paragraph(); continue

        if s.startswith("### "):   _h(doc, s[4:], 3)
        elif s.startswith("## "):  _h(doc, s[3:], 2)
        elif s.startswith("# "):   _h(doc, s[2:], 1)
        elif s.startswith("> "):   _p_md(doc, s[2:], color=D_TEAL)
        elif s.startswith("- ") or s.startswith("* "):
            _bul(doc, s[2:])
        else:
            _p_md(doc, s)  # handles inline **bold** correctly
    flush()


def _start_two_columns(doc: Document):
    """Insert a continuous section break to begin 2-column newspaper layout."""
    p = doc.add_paragraph()
    pPr = p._p.get_or_add_pPr()
    sectPr = OxmlElement("w:sectPr")
    type_el = OxmlElement("w:type"); type_el.set(qn("w:val"), "continuous")
    cols_el = OxmlElement("w:cols")
    cols_el.set(qn("w:num"), "2"); cols_el.set(qn("w:space"), "720")
    sectPr.append(type_el); sectPr.append(cols_el)
    pPr.append(sectPr)

def _end_two_columns(doc: Document):
    """Insert a continuous section break to return to 1-column layout."""
    p = doc.add_paragraph()
    pPr = p._p.get_or_add_pPr()
    sectPr = OxmlElement("w:sectPr")
    type_el = OxmlElement("w:type"); type_el.set(qn("w:val"), "continuous")
    cols_el = OxmlElement("w:cols"); cols_el.set(qn("w:num"), "1")
    sectPr.append(type_el); sectPr.append(cols_el)
    pPr.append(sectPr)

# ─────────────────────────────────────────────────────────────────────────────
# FICHA HELPERS — Student worksheet components
# ─────────────────────────────────────────────────────────────────────────────
def _info_estudiante(doc: Document):
    """Name / Date / Section fields at top of student worksheet."""
    tbl = doc.add_table(rows=1, cols=3)
    tbl.style = "Table Grid"
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    campos = [("Nombre", 3.5), ("Fecha", 1.8), ("Sección", 1.8)]
    for i, (label, w_in) in enumerate(campos):
        cell = tbl.rows[0].cells[i]
        cell.width = Inches(w_in)
        _cell_bg(cell, "F1F8E9")
        _row_h(tbl.rows[0], 0.7)
        p = cell.paragraphs[0]
        r1 = p.add_run(f"{label}: ")
        r1.font.bold = True; r1.font.size = Pt(11)
        r1.font.name = FONT_BODY; r1.font.color.rgb = D_AZUL
        r2 = p.add_run("_" * (int(w_in * 8)))
        r2.font.size = Pt(11); r2.font.name = FONT_BODY
        r2.font.color.rgb = D_GRIS_CLR
    doc.add_paragraph()

def _historia_gancho_box(doc: Document, texto: str):
    """Italic boxed story hook paragraph."""
    if not texto: return
    tbl = doc.add_table(rows=1, cols=1)
    tbl.style = "Table Grid"
    cell = tbl.rows[0].cells[0]
    _cell_bg(cell, "EEF4FF")
    p = cell.paragraphs[0]
    p.paragraph_format.space_before = Pt(4)
    p.paragraph_format.space_after  = Pt(4)
    p.paragraph_format.left_indent  = Inches(0.1)
    run = p.add_run(str(texto))
    run.font.name = FONT_BODY; run.font.size = Pt(11)
    run.font.italic = True; run.font.color.rgb = D_NEGRO
    doc.add_paragraph()

def _opcion_multiple(doc: Document, pregunta: str, opciones: list,
                     letra_colors: list):
    """Multiple choice question with circle markers."""
    # Question
    p_q = doc.add_paragraph()
    p_q.paragraph_format.space_before = Pt(6)
    p_q.paragraph_format.space_after  = Pt(2)
    _runs_from_inline_md(p_q, pregunta, size=12)
    for run in p_q.runs:
        run.font.bold = True

    # Options A) B) C)
    letras = ["A)", "B)", "C)", "D)"]
    for j, opcion in enumerate(opciones):
        p_op = doc.add_paragraph()
        p_op.paragraph_format.left_indent  = Inches(0.4)
        p_op.paragraph_format.space_before = Pt(2)
        p_op.paragraph_format.space_after  = Pt(2)
        r_circ = p_op.add_run("○  ")
        r_circ.font.size = Pt(14); r_circ.font.color.rgb = D_GRIS
        r_letra = p_op.add_run(letras[j] + " ")
        r_letra.font.bold = True; r_letra.font.size = Pt(11)
        r_letra.font.color.rgb = letra_colors[j % len(letra_colors)]
        r_text = p_op.add_run(str(opcion).replace("**","").lstrip("abcdefgABCDEFG) "))
        r_text.font.size = Pt(11); r_text.font.name = FONT_BODY
        r_text.font.color.rgb = D_NEGRO
    doc.add_paragraph()

def _lineas_respuesta(doc: Document, n: int = 3, label: str = ""):
    """Dotted write lines for student responses."""
    if label:
        p_lbl = doc.add_paragraph()
        r = p_lbl.add_run(label)
        r.font.bold = True; r.font.size = Pt(11)
        r.font.color.rgb = D_AZUL
    for _ in range(n):
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(1)
        p.paragraph_format.space_after  = Pt(1)
        run = p.add_run("_" * 90)
        run.font.color.rgb = D_GRIS_CLR; run.font.size = Pt(10)
    doc.add_paragraph()

def _caja_escritura(doc: Document, alto_cm: float = 5.0, instruccion: str = ""):
    """Large bordered box for drawing or extended writing."""
    if instruccion:
        p_i = doc.add_paragraph()
        r = p_i.add_run(instruccion)
        r.font.italic = True; r.font.size = Pt(10)
        r.font.color.rgb = D_GRIS
    tbl = doc.add_table(rows=1, cols=1)
    tbl.style = "Table Grid"
    cell = tbl.rows[0].cells[0]
    _cell_bg(cell, "FFFFFF")
    _row_h(tbl.rows[0], alto_cm)
    p = cell.paragraphs[0]
    p.paragraph_format.space_before = Pt(2)
    r = p.add_run(" ")
    r.font.size = Pt(8)
    doc.add_paragraph()

def _palabra_badge(para, palabra: str, color: RGBColor):
    """Inline word badge for word bank."""
    r1 = para.add_run(f"  ▸ {palabra.upper()}  ")
    r1.font.name = FONT_BODY; r1.font.bold = True
    r1.font.size = Pt(11); r1.font.color.rgb = color

def _banco_palabras(doc: Document, palabras: list):
    """Styled word bank row."""
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(4)
    p.paragraph_format.space_after  = Pt(8)
    colors = [D_AZUL, D_ROJO, D_VERDE2, D_MORADO, D_NARANJA]
    label = p.add_run("Banco de palabras:  ")
    label.font.bold = True; label.font.size = Pt(11); label.font.color.rgb = D_NEGRO
    for i, pal in enumerate(palabras):
        _palabra_badge(p, str(pal), colors[i % len(colors)])

def _generar_sopa_letras(palabras: list, size: int = 14) -> list:
    """
    Word-search: ~1/3 horizontal, ~2/3 vertical, interleaved.
    Prefers positions that cross (overlap) already-placed words.
    """
    grid = [["" for _ in range(size)] for _ in range(size)]
    letras_relleno = "ABCDEFGHIJKLMNOPRSTUVWXZ"

    words = [p.upper().replace(" ", "").replace("Ñ", "N") for p in palabras]
    words = [w for w in words if 0 < len(w) <= size]
    if not words:
        for r in range(size):
            for c in range(size):
                grid[r][c] = random.choice(letras_relleno)
        return grid

    n = len(words)
    n_horiz = max(1, round(n / 3))          # ~1/3 horizontal
    directions = ['H'] * n_horiz + ['V'] * (n - n_horiz)
    random.shuffle(directions)               # interleave so it's not all-H then all-V

    def _find_positions(word, direction):
        """Return (overlap_spots, empty_spots) — prefer overlapping crossings."""
        ln = len(word)
        over, free = [], []
        if direction == 'H':
            for r in range(size):
                for c in range(size - ln + 1):
                    if all(grid[r][c+i] == "" or grid[r][c+i] == word[i]
                           for i in range(ln)):
                        hits = sum(1 for i in range(ln) if grid[r][c+i] == word[i])
                        (over if hits else free).append((r, c))
        else:
            for c in range(size):
                for r in range(size - ln + 1):
                    if all(grid[r+i][c] == "" or grid[r+i][c] == word[i]
                           for i in range(ln)):
                        hits = sum(1 for i in range(ln) if grid[r+i][c] == word[i])
                        (over if hits else free).append((r, c))
        return over, free

    def _place(word, direction, r, c):
        ln = len(word)
        if direction == 'H':
            for i in range(ln): grid[r][c+i] = word[i]
        else:
            for i in range(ln): grid[r+i][c] = word[i]

    for word, direction in zip(words, directions):
        over, free = _find_positions(word, direction)
        pool = over if over else free
        if pool:
            r, c = random.choice(pool)
            _place(word, direction, r, c)
        else:
            # Fallback: try the other direction
            alt = 'V' if direction == 'H' else 'H'
            over2, free2 = _find_positions(word, alt)
            pool2 = over2 if over2 else free2
            if pool2:
                r, c = random.choice(pool2)
                _place(word, alt, r, c)

    for r in range(size):
        for c in range(size):
            if not grid[r][c]:
                grid[r][c] = random.choice(letras_relleno)
    return grid

def _insertar_sopa_docx(doc: Document, grid: list):
    """Render word search grid as a table."""
    size = len(grid)
    tbl = doc.add_table(rows=size, cols=size)
    tbl.style = "Table Grid"
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    cell_w = Inches(0.28)
    for r in range(size):
        row = tbl.rows[r]
        _row_h(row, 0.38)
        for c in range(size):
            cell = row.cells[c]
            cell.width = cell_w
            p = cell.paragraphs[0]; p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run(str(grid[r][c]))
            run.font.name = FONT_BODY; run.font.bold = True
            run.font.size = Pt(10); run.font.color.rgb = D_NEGRO
    doc.add_paragraph()

def _frase_pergamino(doc: Document, frase: str, palabras_secretas: list):
    """Boxed fill-in-the-blank sentence."""
    tbl = doc.add_table(rows=1, cols=1)
    tbl.style = "Table Grid"
    cell = tbl.rows[0].cells[0]
    _cell_bg(cell, "FFF8E1")
    _row_h(tbl.rows[0], 1.4)
    p = cell.paragraphs[0]
    p.paragraph_format.space_before = Pt(6)
    p.paragraph_format.space_after  = Pt(6)
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(str(frase))
    run.font.name = FONT_BODY; run.font.italic = True
    run.font.size = Pt(13); run.font.color.rgb = D_NEGRO
    doc.add_paragraph()
    if palabras_secretas:
        _banco_palabras(doc, palabras_secretas)

# ─────────────────────────────────────────────────────────────────────────────
# DOCX GENERATORS
# ─────────────────────────────────────────────────────────────────────────────
def generar_docx_sintesis(res_m0a: dict, diagnosticos: str) -> bytes:
    doc    = _nuevo_doc()
    unidad = res_m0a.get("unidad_sintetizada", {})
    titulo = unidad.get("titulo", "Síntesis de Unidad")
    _worksheet_header(doc, "Planificación Curricular", titulo)

    _h(doc, "Perfil Neuro-Inclusivo del Aula", 2)
    _p(doc, diagnosticos)
    doc.add_paragraph()

    _h(doc, "Temas de la Unidad", 2)
    for tema in unidad.get("temas_desarrollados", []):
        _h(doc, tema.get("nombre","Tema"), 3, D_VERDE)
        _p(doc, "Conceptos Clave:", bold=True, color=D_AZUL)
        for c in tema.get("conceptos_clave") or tema.get("conceptos") or []:
            _bul(doc, str(c), D_AZUL)
        doc.add_paragraph()
        _p(doc, "Inteligencias Múltiples sugeridas:", bold=True, color=D_VERDE)
        for i in tema.get("inteligencias_sugeridas") or tema.get("inteligencias") or []:
            _bul(doc, str(i))
        doc.add_paragraph()

    _h(doc, "Notas DUA del Docente", 2)
    _p_md(doc, unidad.get("notas_docente","—"))
    doc.add_paragraph()

    _h(doc, "Proyecto ABP — Resumen", 2)
    _p_md(doc, unidad.get("proyecto_pbl","—"))

    buf = io.BytesIO(); doc.save(buf); return buf.getvalue()


def generar_docx_abp(res_m0b: str, res_m0c: str, titulo_unidad: str) -> bytes:
    doc = _nuevo_doc()
    _worksheet_header(doc, "Proyecto ABP", titulo_unidad)
    _h(doc, "Diseño del Proyecto ABP", 2)
    _md(doc, res_m0b or "No generado.")
    doc.add_page_break()
    _h(doc, "Rúbrica y Fichas de Proceso", 2)
    _md(doc, res_m0c or "No generado.")
    buf = io.BytesIO(); doc.save(buf); return buf.getvalue()


def generar_docx_plan_clase(res_m1a: dict, tema: str, diagnosticos: str) -> bytes:
    doc = _nuevo_doc()
    _worksheet_header(doc, "Plan de Clase · 45 minutos", tema)

    verbos = res_m1a.get("mapa_cognitivo", {}).get("verbos", [])
    if verbos:
        _h(doc, "Verbos Operativos — Taxonomía Bloom", 2)
        _bloom_badge_table(doc, verbos)

    im = res_m1a.get("inteligencias_multiples", [])
    if im:
        _h(doc, "Inteligencias Múltiples", 2)
        _tbl(doc, im)

    _section_banner(doc, "⏱️  Secuencia Didáctica — 45 minutos", "1A3A6A", size=13)
    _secuencia_timeline(doc, res_m1a.get("secuencia_didactica", {}).get("bloques", []))

    _section_banner(doc, "🧠  Directrices DUA — Neuro-Inclusión", "007A6E", size=13)
    _dua_framework_header(doc)
    for d in res_m1a.get("dua_neuroinclusion", []):
        dl = str(d).lower()
        if "tea" in dl:
            _dua_note(doc, "👁 TEA", str(d), "007A6E")
        elif "tdah" in dl or "adhd" in dl:
            _dua_note(doc, "⚡ TDAH", str(d), "CC5500")
        else:
            _bul(doc, str(d), size=11)
    doc.add_paragraph()

    adapt = res_m1a.get("tabla_adaptaciones_clase", [])
    if adapt:
        _section_banner(doc, "♿  Adaptaciones por Diagnóstico", "1A5276", size=13)
        _tbl(doc, adapt)

    _h(doc, "Perfil del Aula", 2)
    _p_md(doc, res_m1a.get("perfil_aula_resumido", diagnosticos), size=11)

    buf = io.BytesIO(); doc.save(buf); return buf.getvalue()


def generar_docx_ficha(res_m1c: dict, tema: str) -> bytes:
    """
    Ficha Gamificada — diseño de hoja de trabajo real para imprimir.
    Incluye: espacios de respuesta, sopa de letras real, caja de dibujo.
    """
    doc   = _nuevo_doc()
    raw   = res_m1c
    ficha = raw.get("ficha_trabajo", raw) if isinstance(raw, dict) else {}

    # ── HOJA DEL ESTUDIANTE ──────────────────────────────────────────────────
    _worksheet_header(doc, "Comunicación — Ficha Gamificada", tema)
    _info_estudiante(doc)
    _historia_gancho_box(doc, ficha.get("historia_gancho",""))

    misiones  = ficha.get("misiones", {})
    MC_COLORS = [D_AZUL, D_ROJO, D_VERDE2]

    # ── MISIÓN 1: EL ORÁCULO ─────────────────────────────────────────────────
    _section_banner(doc, "🔮  Misión 1: El Oráculo", "2E6DA4", size=13)
    p_inst = doc.add_paragraph()
    r = p_inst.add_run("Instrucciones: Lee cada pregunta y marca con ○ la respuesta correcta.")
    r.font.italic = True; r.font.size = Pt(10); r.font.color.rgb = D_GRIS
    doc.add_paragraph()

    for idx, q in enumerate(misiones.get("oraculo", [])):
        _opcion_multiple(doc,
                         f"{idx+1}. {q.get('pregunta','')}",
                         q.get("opciones", []),
                         MC_COLORS)

    # ── MISIÓN 2: EL PUENTE ──────────────────────────────────────────────────
    _section_banner(doc, "🌉  Misión 2: El Puente", "007A6E", size=13)
    p_inst2 = doc.add_paragraph()
    r2 = p_inst2.add_run("Instrucciones: Une con una línea cada palabra con su significado.")
    r2.font.italic = True; r2.font.size = Pt(10); r2.font.color.rgb = D_GRIS
    doc.add_paragraph()

    puente_data = misiones.get("puente", [])
    if puente_data:
        # Two-column matching table — words left, definitions right (shuffled)
        import random as _rnd
        definiciones = [p.get("significado","") for p in puente_data]
        definiciones_mezcladas = definiciones[:]
        _rnd.shuffle(definiciones_mezcladas)
        tbl_puente = doc.add_table(rows=len(puente_data)+1, cols=2)
        tbl_puente.style = "Table Grid"
        tbl_puente.alignment = WD_TABLE_ALIGNMENT.CENTER
        # Header
        _cell_bg(tbl_puente.rows[0].cells[0], "007A6E")
        _cell_bg(tbl_puente.rows[0].cells[1], "007A6E")
        for j, hdr_txt in enumerate(["PALABRA", "SIGNIFICADO"]):
            ph = tbl_puente.rows[0].cells[j].paragraphs[0]
            ph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            rh = ph.add_run(hdr_txt)
            rh.font.bold = True; rh.font.color.rgb = D_BLANCO
            rh.font.name = FONT_BODY; rh.font.size = Pt(11)
        # Data rows
        for k, par in enumerate(puente_data):
            row = tbl_puente.rows[k+1]
            _cell_bg(row.cells[0], "E8F5E9")
            _cell_bg(row.cells[1], "FFFFFF")
            _row_h(row, 0.7)
            p_pal = row.cells[0].paragraphs[0]
            p_pal.alignment = WD_ALIGN_PARAGRAPH.CENTER
            r_pal = p_pal.add_run(str(par.get("palabra","")))
            r_pal.font.bold = True; r_pal.font.color.rgb = D_VERDE
            r_pal.font.name = FONT_BODY; r_pal.font.size = Pt(11)
            p_def = row.cells[1].paragraphs[0]
            r_def = p_def.add_run(str(definiciones_mezcladas[k]))
            r_def.font.color.rgb = D_NEGRO
            r_def.font.name = FONT_BODY; r_def.font.size = Pt(10)
    doc.add_paragraph()

    # ── MISIÓN 3: LA SOPA DE LETRAS ─────────────────────────────────────────
    _section_banner(doc, "🥣  Misión 3: La Sopa de Letras", "1A5276", size=13)
    p_inst3 = doc.add_paragraph()
    r3 = p_inst3.add_run("Instrucciones: Encuentra y encierra las siguientes palabras en la sopa de letras.")
    r3.font.italic = True; r3.font.size = Pt(10); r3.font.color.rgb = D_GRIS
    doc.add_paragraph()

    palabras_sopa = misiones.get("sopa", [])
    if palabras_sopa:
        _banco_palabras(doc, palabras_sopa)
        grid = _generar_sopa_letras(palabras_sopa, size=12)
        _insertar_sopa_docx(doc, grid)

    # ── MISIÓN 4: EL PERGAMINO ───────────────────────────────────────────────
    _section_banner(doc, "📜  Misión 4: El Pergamino", "6A3D9A", size=13)
    p_inst4 = doc.add_paragraph()
    r4 = p_inst4.add_run("Instrucciones: Completa la frase usando las palabras del banco.")
    r4.font.italic = True; r4.font.size = Pt(10); r4.font.color.rgb = D_GRIS
    doc.add_paragraph()

    per = misiones.get("pergamino", {})
    _frase_pergamino(doc, per.get("frase_con_espacios",""), per.get("palabras_secretas",[]))

    # ── MISIÓN 5: EL LIENZO ──────────────────────────────────────────────────
    _section_banner(doc, "🎨  Misión 5: El Lienzo", "CC5500", size=13)
    lienzo_inst = str(misiones.get("lienzo","Crea tu obra de arte basándote en lo aprendido."))
    p_inst5 = doc.add_paragraph()
    r5 = p_inst5.add_run(lienzo_inst.replace("**",""))
    r5.font.size = Pt(11); r5.font.color.rgb = D_NEGRO; r5.font.name = FONT_BODY
    doc.add_paragraph()
    _caja_escritura(doc, alto_cm=6.0,
                    instruccion="✏️ Dibuja o escribe tu respuesta en este espacio:")

    # ── PÁGINA DE ADAPTACIONES (solo para el docente) ────────────────────────
    adapt = ficha.get("adaptaciones_por_mision", [])
    if adapt:
        doc.add_page_break()
        _section_banner(doc, "📋  ANEXO DOCENTE — Adaptaciones por Misión", "1A3A6A", size=13)
        p_nota = doc.add_paragraph()
        r_nota = p_nota.add_run("Esta página es solo para el docente. No imprimir para los estudiantes.")
        r_nota.font.italic = True; r_nota.font.color.rgb = D_ROJO; r_nota.font.size = Pt(10)
        doc.add_paragraph()
        _tbl(doc, adapt)

    buf = io.BytesIO(); doc.save(buf); return buf.getvalue()


def generar_docx_evaluaciones(res_m2a: str, res_m2b: str, tema: str) -> bytes:
    doc = _nuevo_doc()

    # ── HOJA DEL ESTUDIANTE: Pop Quiz ────────────────────────────────────────
    _worksheet_header(doc, "Evaluación Formativa DUA", tema)
    _info_estudiante(doc)
    _section_banner(doc, "🎲  Pop Quiz — ¡Misión de Cierre!", "2E6DA4", size=13)
    _md(doc, res_m2a or "No generado.")

    # ── PÁGINA DEL DOCENTE: Panel de Control ─────────────────────────────────
    if res_m2b:
        doc.add_page_break()
        _section_banner(doc, "📊  PANEL DE CONTROL DEL TUTOR  —  Solo para el Docente",
                        "1A5276", size=13)
        p_nota = doc.add_paragraph()
        r_nota = p_nota.add_run("Esta página es confidencial. No distribuir a los estudiantes.")
        r_nota.font.italic = True; r_nota.font.color.rgb = D_ROJO; r_nota.font.size = Pt(10)
        doc.add_paragraph()
        _start_two_columns(doc)
        _md(doc, res_m2b)
        _end_two_columns(doc)

    buf = io.BytesIO(); doc.save(buf); return buf.getvalue()


# ─────────────────────────────────────────────────────────────────────────────
# HTML GENERATORS
# ─────────────────────────────────────────────────────────────────────────────
def _logo_b64() -> str:
    """Return logo as base64 data URI, or empty string if not found."""
    import base64
    if os.path.exists(LOGO_PATH):
        with open(LOGO_PATH, "rb") as f:
            return "data:image/png;base64," + base64.b64encode(f.read()).decode()
    return ""


def _md_to_html(text: str) -> str:
    """Minimal markdown → HTML: headings, bullets, bold, paragraphs."""
    if not text:
        return ""
    lines = str(text).split("\n")
    out   = []
    in_ul = False

    def close_ul():
        nonlocal in_ul
        if in_ul:
            out.append("</ul>"); in_ul = False

    def inline(s):
        s = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', s)
        s = re.sub(r'\*(.+?)\*',     r'<em>\1</em>', s)
        return s

    for line in lines:
        s = line.strip()
        if not s:
            close_ul(); out.append("<br>"); continue
        if s in ('---', '___', '***') or (len(s) >= 3 and len(set(s)) == 1):
            close_ul(); continue
        if s.startswith("### "): close_ul(); out.append(f"<h3>{inline(s[4:])}</h3>"); continue
        if s.startswith("## "):  close_ul(); out.append(f"<h2>{inline(s[3:])}</h2>"); continue
        if s.startswith("# "):   close_ul(); out.append(f"<h1>{inline(s[2:])}</h1>"); continue
        if s.startswith("- ") or s.startswith("* "):
            if not in_ul: out.append("<ul>"); in_ul = True
            out.append(f"<li>{inline(s[2:])}</li>")
            continue
        close_ul()
        out.append(f"<p>{inline(s)}</p>")

    close_ul()
    return "\n".join(out)


# ─────────────────────────────────────────────────────────────────────────────
# WORKSHEET THEMES
# ─────────────────────────────────────────────────────────────────────────────
# CSS uses normal single braces (not f-string embedded) — inject via {theme_css}

_CSS_SHARED_PRINT = """
@media print{
  @page{size:A4;margin:15mm 12mm 15mm 12mm}
  body{padding:0;font-size:10pt}
  .print-color{-webkit-print-color-adjust:exact;print-color-adjust:exact}
  .teacher-page{page-break-before:always}
  .section,.match-container,.draw-box,.fill-box,.sopa-grid,.word-bank{page-break-inside:avoid;break-inside:avoid}
  a{text-decoration:none}
}
"""

# ── Aventura: bold color, Fredoka One, gamified ─────────────────────────────
_CSS_AVENTURA = """
*{box-sizing:border-box;margin:0;padding:0}
body{font-family:'Nunito',sans-serif;font-size:12pt;color:#1a1a1a;background:#fff;padding:20px}
.header{display:flex;align-items:center;gap:16px;background:#1a3a6a;color:#fff;padding:12px 20px;border-radius:8px;margin-bottom:18px}
.logo{height:56px;width:auto}
.header-text{flex:1}
.header-title{font-family:'Fredoka One',sans-serif;font-size:16pt;color:#fff}
.header-school{font-size:9pt;color:#c8e6c9;margin-top:2px}
.header-grado{font-family:'Fredoka One',sans-serif;font-size:18pt;color:#ffe082;white-space:nowrap}
.tema-title{text-align:center;font-size:17pt;font-weight:800;color:#cc2200;text-decoration:underline;margin:10px 0 16px}
.student-fields{display:flex;gap:12px;margin-bottom:14px}
.field{border:1.5px solid #90a4ae;border-radius:6px;padding:6px 12px;flex:1;font-size:11pt}
.field span{font-weight:700;color:#1a3a6a}
.field .line{border-bottom:1.5px dashed #b0bec5;display:inline-block;width:60%;margin-left:4px}
.points-bar{display:flex;align-items:center;gap:8px;background:#fff8e1;border:2px solid #ffe082;border-radius:8px;padding:8px 14px;margin-bottom:14px;font-size:11pt;font-weight:700;color:#cc5500}
.section{margin-bottom:18px}
.section-header{display:flex;align-items:center;padding:8px 14px;border-radius:6px;color:#fff;font-family:'Fredoka One',sans-serif;font-size:13pt;margin-bottom:10px}
.section-header.blue{background:#2e6da4}.section-header.teal{background:#007a6e}
.section-header.navy{background:#1a5276}.section-header.purple{background:#6a3d9a}
.section-header.orange{background:#cc5500}
.inst{font-style:italic;color:#555;font-size:10pt;margin-bottom:8px}
.gancho{background:#eef4ff;border-left:4px solid #1a3a6a;padding:10px 14px;border-radius:0 6px 6px 0;font-style:italic;margin-bottom:14px}
.mc-question{font-weight:700;font-size:11pt;margin:10px 0 4px}
.mc-options{padding-left:20px;margin-bottom:10px}
.mc-option{margin:3px 0;font-size:11pt}
.mc-option::before{content:"○  ";color:#90a4ae;font-size:13pt}
.match-table{width:100%;border-collapse:collapse;margin-bottom:12px}
.match-table th{background:#007a6e;color:#fff;padding:6px 10px;text-align:center;font-size:10pt}
.match-table td{border:1px solid #b2dfdb;padding:6px 10px;font-size:10pt;height:32px}
.match-table tr:nth-child(even) td{background:#f1f8f6}
.match-table td.word{background:#e8f5e9;font-weight:700;color:#1b5e20;text-align:center}
.match-container{display:flex;justify-content:space-between;margin:14px 0;padding:0 8px;gap:28px}
.match-col{display:flex;flex-direction:column;gap:14px;width:45%}
.match-item{display:flex;align-items:center;background:#e8f5e9;border:1.5px solid #b2dfdb;padding:8px 14px;border-radius:8px;position:relative;font-size:10pt;overflow:visible}
.match-item.left{justify-content:flex-end;text-align:right;padding-right:22px}
.match-item.right{justify-content:flex-start;text-align:left;padding-left:22px}
.node{width:12px;height:12px;background:#007a6e;border-radius:50%;position:absolute;top:50%;transform:translateY(-50%)}
.match-item.left .node{right:-7px}
.match-item.right .node{left:-7px}
.sopa-grid{display:inline-grid;gap:2px;background:#e8f5e9;padding:8px;border-radius:6px;border:1.5px solid #2e7d32;margin:8px 0}
.sopa-grid span{display:flex;align-items:center;justify-content:center;width:24px;height:24px;font-weight:700;font-size:10pt;background:#fff;border-radius:3px}
.word-bank{display:flex;flex-wrap:wrap;gap:8px;margin:8px 0 12px}
.word-chip{background:#1a3a6a;color:#fff;padding:3px 10px;border-radius:20px;font-weight:700;font-size:10pt}
.fill-box{background:#fff8e1;border:1.5px dashed #f9a825;border-radius:6px;padding:10px 14px;font-style:italic;font-size:12pt;text-align:center;margin:8px 0 12px;line-height:2.5}
.fill-blank{display:inline-block;border-bottom:2px solid #1a3a6a;min-width:80px;margin:0 3px}
.answer-lines{margin:8px 0 12px}
.answer-line{border-bottom:1.5px dotted #b0bec5;height:30px;margin-bottom:2px}
.draw-box{border:2px dashed #90a4ae;border-radius:8px;min-height:160px;margin:8px 0 12px;background:#fafafa;display:flex;align-items:flex-end;padding:6px}
.draw-label{color:#90a4ae;font-size:9pt;font-style:italic}
.sent-row{display:flex;gap:20px;align-items:center;margin:10px 0 14px;flex-wrap:wrap}
.sent-item{display:flex;align-items:center;gap:8px;font-size:11pt}
.sent-circle{width:32px;height:32px;border-radius:50%;border:2px solid #1a3a6a;display:inline-flex;align-items:center;justify-content:center;font-size:16pt}
.teacher-page{page-break-before:always;margin-top:40px}
.teacher-banner{background:#1a5276;color:#fff;padding:10px 14px;border-radius:6px;font-family:'Fredoka One',sans-serif;font-size:13pt;margin-bottom:10px}
.confidential{color:#cc2200;font-style:italic;font-size:10pt;margin-bottom:12px}
p{margin:4px 0} ul{margin:4px 0 8px 20px} li{margin:2px 0}
h1{font-size:15pt;color:#1a3a6a} h2{font-size:13pt;color:#1a3a6a} h3{font-size:11pt;color:#1a3a6a}
"""

# ── Cuaderno: warm cream, handwritten headers, naturalist notebook feel ──────
_CSS_CUADERNO = """
*{box-sizing:border-box;margin:0;padding:0}
body{font-family:'Nunito',sans-serif;font-size:12pt;color:#3e2723;background:#fdf8f0;padding:20px}
.header{display:flex;align-items:center;gap:16px;background:#33691e;color:#fff;padding:12px 20px;border-radius:4px;margin-bottom:18px;border-bottom:5px solid #8bc34a}
.logo{height:56px;width:auto}
.header-text{flex:1}
.header-title{font-family:'Caveat',cursive;font-size:21pt;color:#fff;letter-spacing:0.5px}
.header-school{font-size:9pt;color:#dcedc8;margin-top:2px}
.header-grado{font-family:'Caveat',cursive;font-size:24pt;color:#f9a825;white-space:nowrap}
.tema-title{text-align:center;font-size:20pt;font-weight:700;color:#4e342e;font-family:'Caveat',cursive;border-bottom:3px double #a1887f;padding-bottom:6px;margin:10px 0 16px}
.student-fields{display:flex;gap:12px;margin-bottom:14px}
.field{border:1.5px solid #bcaaa4;border-radius:4px;padding:6px 12px;flex:1;font-size:11pt;background:#fffdf7}
.field span{font-weight:700;color:#4e342e}
.field .line{border-bottom:1.5px dashed #bcaaa4;display:inline-block;width:60%;margin-left:4px}
.points-bar{display:flex;align-items:center;gap:8px;background:#fff9c4;border:2px solid #f9a825;border-radius:4px;padding:8px 14px;margin-bottom:14px;font-size:11pt;font-weight:700;color:#e65100}
.section{margin-bottom:18px}
.section-header{display:flex;align-items:center;padding:8px 16px;border-radius:4px;color:#fff;font-family:'Caveat',cursive;font-size:17pt;margin-bottom:10px;letter-spacing:0.3px}
.section-header.blue{background:#558b2f}.section-header.teal{background:#00695c}
.section-header.navy{background:#5d4037}.section-header.purple{background:#6a1e4a}
.section-header.orange{background:#bf360c}
.inst{font-style:italic;color:#795548;font-size:10pt;margin-bottom:8px}
.gancho{background:#fff9c4;border-left:5px solid #f9a825;padding:10px 14px;border-radius:0 4px 4px 0;font-style:italic;margin-bottom:14px;color:#4e342e}
.mc-question{font-weight:700;font-size:11pt;margin:10px 0 4px;color:#3e2723}
.mc-options{padding-left:20px;margin-bottom:10px}
.mc-option{margin:4px 0;font-size:11pt}
.mc-option::before{content:"○  ";color:#a1887f;font-size:13pt}
.match-table{width:100%;border-collapse:collapse;margin-bottom:12px}
.match-table th{background:#558b2f;color:#fff;padding:6px 10px;text-align:center;font-family:'Caveat',cursive;font-size:13pt}
.match-table td{border:1px solid #d7ccc8;padding:6px 10px;font-size:10pt;height:34px;background:#fffdf7}
.match-table tr:nth-child(even) td{background:#f5f0e8}
.match-table td.word{background:#dcedc8;font-weight:700;color:#33691e;text-align:center}
.match-container{display:flex;justify-content:space-between;margin:14px 0;padding:0 8px;gap:28px}
.match-col{display:flex;flex-direction:column;gap:14px;width:45%}
.match-item{display:flex;align-items:center;background:#fffdf7;border:1.5px solid #d7ccc8;padding:8px 14px;border-radius:4px;position:relative;font-size:10pt;overflow:visible}
.match-item.left{justify-content:flex-end;text-align:right;padding-right:22px}
.match-item.right{justify-content:flex-start;text-align:left;padding-left:22px}
.node{width:12px;height:12px;background:#33691e;border-radius:50%;position:absolute;top:50%;transform:translateY(-50%)}
.match-item.left .node{right:-7px}
.match-item.right .node{left:-7px}
.sopa-grid{display:inline-grid;gap:2px;background:#dcedc8;padding:8px;border-radius:4px;border:1.5px solid #558b2f;margin:8px 0}
.sopa-grid span{display:flex;align-items:center;justify-content:center;width:24px;height:24px;font-weight:700;font-size:10pt;background:#fffdf7;border-radius:2px}
.word-bank{display:flex;flex-wrap:wrap;gap:8px;margin:8px 0 12px}
.word-chip{background:#558b2f;color:#fff;padding:3px 10px;border-radius:4px;font-weight:700;font-family:'Caveat',cursive;font-size:13pt}
.fill-box{background:#fffdf7;border:2px dashed #bcaaa4;border-radius:4px;padding:10px 14px;font-style:italic;font-size:12pt;text-align:center;margin:8px 0 12px;line-height:2.8}
.fill-blank{display:inline-block;border-bottom:2px solid #5d4037;min-width:80px;margin:0 3px}
.answer-lines{margin:8px 0 12px}
.answer-line{border-bottom:2px solid #d7ccc8;height:36px;margin-bottom:0}
.draw-box{border:2px dashed #bcaaa4;border-radius:4px;min-height:168px;margin:8px 0 12px;background:#fffdf7;display:flex;align-items:flex-end;padding:6px}
.draw-label{color:#a1887f;font-size:9pt;font-style:italic}
.sent-row{display:flex;gap:20px;align-items:center;margin:10px 0 14px;flex-wrap:wrap}
.sent-item{display:flex;align-items:center;gap:8px;font-size:11pt}
.sent-circle{width:32px;height:32px;border-radius:50%;border:2px solid #558b2f;display:inline-flex;align-items:center;justify-content:center;font-size:16pt}
.teacher-page{page-break-before:always;margin-top:40px}
.teacher-banner{background:#5d4037;color:#fff;padding:10px 14px;border-radius:4px;font-family:'Caveat',cursive;font-size:17pt;margin-bottom:10px}
.confidential{color:#bf360c;font-style:italic;font-size:10pt;margin-bottom:12px}
p{margin:4px 0} ul{margin:4px 0 8px 20px} li{margin:2px 0}
h1{font-size:15pt;color:#4e342e} h2{font-size:13pt;color:#4e342e} h3{font-size:11pt;color:#5d4037}
"""

# ── Periódico: black + crimson editorial, Playfair Display serif, sharp ──────
_CSS_PERIODICO = """
*{box-sizing:border-box;margin:0;padding:0}
body{font-family:'Source Serif 4','Georgia',serif;font-size:12pt;color:#111;background:#fff;padding:20px}
.header{display:flex;align-items:center;gap:16px;background:#111;color:#fff;padding:14px 20px;border-radius:0;margin-bottom:0;border-top:6px solid #c62828;border-bottom:1px solid #333}
.logo{height:56px;width:auto;filter:brightness(0) invert(1)}
.header-text{flex:1;border-left:1px solid #444;padding-left:16px}
.header-title{font-family:'Playfair Display',serif;font-size:15pt;color:#fff;letter-spacing:0.5px}
.header-school{font-size:9pt;color:#999;margin-top:2px}
.header-grado{font-family:'Playfair Display',serif;font-size:16pt;color:#c62828;white-space:nowrap;text-align:right}
.tema-title{text-align:center;font-size:22pt;font-weight:900;color:#111;font-family:'Playfair Display',serif;border-top:3px solid #111;border-bottom:3px solid #111;padding:8px 0;margin:14px 0 18px;letter-spacing:0.5px}
.student-fields{display:flex;gap:12px;margin-bottom:14px}
.field{border:1.5px solid #111;border-radius:0;padding:6px 12px;flex:1;font-size:11pt}
.field span{font-weight:700;color:#111;font-family:'Playfair Display',serif}
.field .line{border-bottom:1.5px solid #888;display:inline-block;width:60%;margin-left:4px}
.points-bar{display:flex;align-items:center;gap:8px;background:#c62828;border:none;border-radius:0;padding:8px 14px;margin-bottom:16px;font-size:11pt;font-weight:700;color:#fff;font-family:'Playfair Display',serif}
.section{margin-bottom:20px;padding-bottom:16px;border-bottom:1px solid #ddd}
.section-header{display:flex;align-items:center;padding:8px 14px;border-radius:0;color:#fff;font-family:'Playfair Display',serif;font-size:14pt;margin-bottom:10px;border-left:6px solid #c62828}
.section-header.blue{background:#111}.section-header.teal{background:#111}
.section-header.navy{background:#111}.section-header.purple{background:#111}
.section-header.orange{background:#111}
.inst{font-style:italic;color:#555;font-size:10pt;margin-bottom:8px}
.gancho{background:#fff;border:2px solid #111;padding:16px 20px;font-style:italic;margin-bottom:14px;font-family:'Playfair Display',serif;font-size:13pt;color:#111;text-align:center;border-top:4px solid #c62828}
.mc-question{font-weight:700;font-size:11pt;margin:10px 0 4px;font-family:'Playfair Display',serif}
.mc-options{padding-left:20px;margin-bottom:10px}
.mc-option{margin:4px 0;font-size:11pt}
.mc-option::before{content:"○  ";color:#555;font-size:13pt}
.match-table{width:100%;border-collapse:collapse;margin-bottom:12px;border:2px solid #111}
.match-table th{background:#111;color:#fff;padding:6px 10px;text-align:center;font-family:'Playfair Display',serif;font-size:11pt;border:1px solid #333}
.match-table td{border:1px solid #888;padding:6px 10px;font-size:10pt;height:32px}
.match-table tr:nth-child(even) td{background:#f5f5f5}
.match-table td.word{background:#c62828;color:#fff;font-weight:700;text-align:center;border-color:#c62828}
.match-container{display:flex;justify-content:space-between;margin:14px 0;padding:0 8px;gap:28px}
.match-col{display:flex;flex-direction:column;gap:14px;width:45%}
.match-item{display:flex;align-items:center;background:#f9f9f9;border:1.5px solid #999;padding:8px 14px;border-radius:0;position:relative;font-size:10pt;overflow:visible}
.match-item.left{justify-content:flex-end;text-align:right;padding-right:22px}
.match-item.right{justify-content:flex-start;text-align:left;padding-left:22px}
.node{width:12px;height:12px;background:#111;border-radius:0;position:absolute;top:50%;transform:translateY(-50%)}
.match-item.left .node{right:-7px}
.match-item.right .node{left:-7px}
.sopa-grid{display:inline-grid;gap:1px;background:#111;padding:4px;border-radius:0;border:2px solid #111;margin:8px 0}
.sopa-grid span{display:flex;align-items:center;justify-content:center;width:26px;height:26px;font-weight:700;font-size:10pt;background:#fff;border-radius:0;font-family:'Courier New',monospace}
.word-bank{display:flex;flex-wrap:wrap;gap:8px;margin:8px 0 12px}
.word-chip{background:#111;color:#fff;padding:3px 10px;border-radius:0;font-weight:700;font-size:10pt;font-family:'Playfair Display',serif}
.fill-box{background:#f5f5f5;border:2px solid #111;border-radius:0;padding:10px 14px;font-style:italic;font-size:12pt;text-align:center;margin:8px 0 12px;line-height:2.6}
.fill-blank{display:inline-block;border-bottom:2px solid #111;min-width:90px;margin:0 3px}
.answer-lines{margin:8px 0 12px}
.answer-line{border-bottom:1.5px solid #bbb;height:32px;margin-bottom:2px}
.draw-box{border:2px solid #111;border-radius:0;min-height:160px;margin:8px 0 12px;background:#fafafa;display:flex;align-items:flex-end;padding:6px}
.draw-label{color:#888;font-size:9pt;font-style:italic}
.sent-row{display:flex;gap:20px;align-items:center;margin:10px 0 14px;flex-wrap:wrap}
.sent-item{display:flex;align-items:center;gap:8px;font-size:11pt}
.sent-circle{width:32px;height:32px;border-radius:0;border:2px solid #111;display:inline-flex;align-items:center;justify-content:center;font-size:16pt}
.teacher-page{page-break-before:always;margin-top:40px}
.teacher-banner{background:#c62828;color:#fff;padding:10px 14px;border-radius:0;font-family:'Playfair Display',serif;font-size:14pt;margin-bottom:10px;letter-spacing:0.3px}
.confidential{color:#c62828;font-style:italic;font-size:10pt;margin-bottom:12px}
p{margin:4px 0} ul{margin:4px 0 8px 20px} li{margin:2px 0}
h1{font-size:15pt;color:#111;font-family:'Playfair Display',serif}
h2{font-size:13pt;color:#111;font-family:'Playfair Display',serif}
h3{font-size:11pt;color:#111;font-family:'Playfair Display',serif}
"""

WORKSHEET_THEMES = {
    "aventura": {
        "label": "🏆 Aventura — colorido, gamificado",
        "gfonts": "family=Nunito:wght@400;600;700;800&family=Fredoka+One",
        "css": _CSS_AVENTURA + _CSS_SHARED_PRINT,
    },
    "cuaderno": {
        "label": "📓 Cuaderno de Campo — cálido, cuaderno de naturalista",
        "gfonts": "family=Nunito:wght@400;600;700;800&family=Caveat:wght@600;700",
        "css": _CSS_CUADERNO + _CSS_SHARED_PRINT,
    },
    "periodico": {
        "label": "📰 El Periódico — editorial, tipografía serif, alto contraste",
        "gfonts": "family=Playfair+Display:wght@700;900&family=Source+Serif+4:wght@400;600",
        "css": _CSS_PERIODICO + _CSS_SHARED_PRINT,
    },
}


def _render_html_doc(title: str, materia: str, tema: str,
                     body_html: str, grado: str = "5to",
                     theme: str = "aventura") -> str:
    """Base Las Palmas HTML document template — print-ready A4."""
    tc       = WORKSHEET_THEMES.get(theme, WORKSHEET_THEMES["aventura"])
    logo     = _logo_b64()
    logo_tag = f'<img src="{logo}" alt="Las Palmas School" class="logo">' if logo else ""
    gfonts   = tc["gfonts"]
    css      = tc["css"]
    return f"""<!DOCTYPE html>
<html lang="es">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1.0">
<title>{title}</title>
<link href="https://fonts.googleapis.com/css2?{gfonts}&display=swap" rel="stylesheet">
<style>{css}</style>
</head>
<body>
<div class="header print-color">
  {logo_tag}
  <div class="header-text">
    <div class="header-title">Fichas de Trabajo · Las Palmas School</div>
    <div class="header-school">{materia}</div>
  </div>
  <div class="header-grado">{grado}<br>PRIMARIA</div>
</div>
<div class="tema-title">{tema.upper()}</div>
{body_html}
</body>
</html>"""


def _sopa_html(palabras: list, size: int = 14) -> str:
    """Generate word-search HTML as a plain <table> (wkhtmltopdf-compatible)."""
    grid = _generar_sopa_letras(palabras, size)
    TD = ('style="width:28px;height:28px;text-align:center;vertical-align:middle;'
          'font-family:monospace;font-weight:bold;font-size:13px;'
          'border:1px solid #ccc;color:#1a1a1a;"')
    rows = "".join(
        "<tr>" + "".join(f"<td {TD}>{grid[r][c]}</td>" for c in range(size)) + "</tr>"
        for r in range(size)
    )
    return f'<table style="border-collapse:collapse;margin:12px auto;">{rows}</table>'


_INTRO_IA_EVAL = [
    "soy el motor m2", "método palma-ribera", "he diseñado este",
    "el objetivo es evaluar de forma rápida", "basado en la planificación estratégica",
    "hola! soy", "¡hola! soy el motor",
]

def _es_intro_eval(linea: str) -> bool:
    l = linea.lower().strip()
    return any(m in l for m in _INTRO_IA_EVAL)

def _limpiar_eval(txt: str) -> str:
    salida = []
    for l in txt.split("\n"):
        if _es_intro_eval(l): continue
        s = l.strip()
        if s and len(set(s)) == 1 and s[0] in "-_*": continue
        if re.fullmatch(r"[#*]+\s*", s): continue
        salida.append(l)
    return "\n".join(salida)

def _parsear_secciones_eval(texto: str) -> dict:
    secciones = {"intro": [], "m1": [], "m2": [], "m3": [], "antes": []}
    actual = "intro"
    PATRONES = {
        "m1":    re.compile(r"mis[ií]on\s*1|el artista|dibujo", re.I),
        "m2":    re.compile(r"mis[ií]on\s*2|el historiador|escritura\s+corta", re.I),
        "m3":    re.compile(r"mis[ií]on\s*3|el viajero|oral|m[ií]mica", re.I),
        "antes": re.compile(r"antes de irte|c[oó]mo te sentiste", re.I),
    }
    for linea in texto.split("\n"):
        s_raw = linea.strip()
        for key, pat in PATRONES.items():
            if pat.search(s_raw):
                actual = key
                break
        secciones[actual].append(linea)
    return {k: "\n".join(v).strip() for k, v in secciones.items()}

def _strip_hdr_eval(txt: str) -> str:
    lines = txt.split("\n")
    skip = re.compile(r"^#|mis[ií]on\s*[123]|artista|historiador|viajero", re.I)
    while lines and skip.search(lines[0].strip()):
        lines = lines[1:]
    return "\n".join(lines).strip()

def _sentiment_html() -> str:
    items = [("😎", "¡Entendí todo!"), ("🤔", "Tengo algunas dudas"), ("🤯", "¡Fue mucha información!")]
    inner = "".join(
        f'<div class="sent-item"><span class="sent-circle">{e}</span><span>{t}</span></div>'
        for e, t in items
    )
    return f'<div class="sent-row">{inner}</div>'


def generar_html_evaluaciones(res_m2a: str, res_m2b: str, tema: str,
                              theme: str = "aventura") -> str:
    """
    Ticket de Salida + Guía Docente como HTML listo para imprimir.
    Limpia intro IA, parsea secciones M1/M2/M3, agrega espacios reales por tipo.
    """
    texto = _limpiar_eval(res_m2a or "")
    secs  = _parsear_secciones_eval(texto)

    body = []
    body.append("""<div class="student-fields">
  <div class="field"><span>Nombre:</span> <span class="line"></span></div>
  <div class="field"><span>Fecha:</span> <span class="line"></span></div>
  <div class="field"><span>Sección:</span> <span class="line"></span></div>
</div>""")
    body.append('<div class="points-bar">⭐ Puntuación: ___ / ___</div>')
    body.append('<div class="section">')
    body.append('<div class="section-header blue">🎲 Pop Quiz — ¡Misión de Cierre!</div>')

    if secs["intro"].strip():
        body.append(_md_to_html(secs["intro"]))

    if secs["m1"].strip():
        body.append('<div class="section-header teal">🎨 Misión 1: El Artista (Dibujo)</div>')
        body.append(_md_to_html(_strip_hdr_eval(secs["m1"])))
        body.append('<div class="draw-box"><span class="draw-label">✏️ Dibuja tu respuesta aquí</span></div>')

    if secs["m2"].strip():
        body.append('<div class="section-header navy">✍️ Misión 2: El Historiador (Escritura corta)</div>')
        body.append(_md_to_html(_strip_hdr_eval(secs["m2"])))
        body.append('<div class="answer-lines">' + '<div class="answer-line"></div>' * 4 + '</div>')

    if secs["m3"].strip():
        body.append('<div class="section-header purple">🗣️ Misión 3: El Viajero (Oral o Mímica)</div>')
        body.append(_md_to_html(_strip_hdr_eval(secs["m3"])))

    if secs["antes"].strip():
        texto_antes = re.sub(r"[😎🤔🤯].*", "", secs["antes"]).strip()
        texto_antes = re.sub(r"\*\*", "", texto_antes)
        body.append(f'<p style="margin-top:12px"><strong>Antes de irte:</strong> {texto_antes}</p>')
        body.append(_sentiment_html())

    body.append('</div>')

    teacher_html = ""
    if res_m2b and res_m2b.strip():
        lineas_b = [l for l in res_m2b.split("\n") if not _es_intro_eval(l)]
        teacher_html = f"""
<div class="teacher-page">
  <div class="teacher-banner">📋 CLAVE DE RESPUESTAS Y GUÍA DOCENTE — No imprimir para los estudiantes</div>
  <p class="confidential">Uso exclusivo del docente.</p>
  {_md_to_html(chr(10).join(lineas_b))}
</div>"""

    return _render_html_doc(
        title    = f"Evaluación — {tema}",
        materia  = "Evaluación Formativa DUA",
        tema     = tema,
        body_html = "\n".join(body) + teacher_html,
        theme    = theme,
    )


def generar_html_ficha(res_m1c: dict, tema: str, theme: str = "aventura") -> str:
    """
    Ficha Gamificada as print-ready HTML (Google Fonts, full-color sections,
    CSS word-search grid, print CSS, teacher adaptations on separate page).
    """
    raw   = res_m1c
    ficha = raw.get("ficha_trabajo", raw) if isinstance(raw, dict) else {}
    misiones = ficha.get("misiones", {})
    body     = []

    # ── Student info fields ──────────────────────────────────────────────────
    body.append("""<div class="student-fields">
  <div class="field"><span>Nombre:</span> <span class="line"></span></div>
  <div class="field"><span>Fecha:</span> <span class="line"></span></div>
  <div class="field"><span>Sección:</span> <span class="line"></span></div>
</div>""")

    # ── Points tracker ───────────────────────────────────────────────────────
    body.append('<div class="points-bar">⭐ Colecciona tus puntos: ___ / 5</div>')

    # ── Hook story ───────────────────────────────────────────────────────────
    gancho = str(ficha.get("historia_gancho","")).strip()
    if gancho:
        body.append(f'<div class="gancho">{gancho}</div>')

    # ── Misión 1: El Oráculo (MC) ────────────────────────────────────────────
    body.append('<div class="section">')
    body.append('<div class="section-header blue">🔮 Misión 1: El Oráculo</div>')
    body.append('<p class="inst">Lee cada pregunta y marca con ○ la respuesta correcta.</p>')
    for q in misiones.get("oraculo", []):
        pregunta = str(q.get("pregunta",""))
        opts     = q.get("opciones", [])
        body.append(f'<div class="mc-question">{pregunta}</div>')
        body.append('<div class="mc-options">')
        letras = ["A)", "B)", "C)", "D)"]
        for j, opt in enumerate(opts):
            clean_opt = str(opt).lstrip("abcdefgABCDEFG) ").replace("**","")
            body.append(f'<div class="mc-option">{letras[j]} {clean_opt}</div>')
        body.append('</div>')
    body.append('</div>')

    # ── Misión 2: El Puente (matching) ───────────────────────────────────────
    puente = misiones.get("puente", [])
    if puente:
        import random as _rnd
        defs_mezcladas = [p.get("significado","") for p in puente]
        _rnd.shuffle(defs_mezcladas)
        body.append('<div class="section">')
        body.append('<div class="section-header teal">🌉 Misión 2: El Puente</div>')
        body.append('<p class="inst">Une con una línea cada palabra con su significado.</p>')
        left_items  = "".join(
            f'<div class="match-item left">{str(par.get("palabra",""))}<span class="node"></span></div>'
            for par in puente
        )
        right_items = "".join(
            f'<div class="match-item right"><span class="node"></span>{defs_mezcladas[k] if k < len(defs_mezcladas) else ""}</div>'
            for k in range(len(puente))
        )
        body.append(
            f'<div class="match-container">'
            f'<div class="match-col">{left_items}</div>'
            f'<div class="match-col">{right_items}</div>'
            f'</div>'
        )
        body.append('</div>')

    # ── Misión 3: Sopa de Letras ─────────────────────────────────────────────
    palabras_sopa = misiones.get("sopa", [])
    if palabras_sopa:
        body.append('<div class="section">')
        body.append('<div class="section-header navy">🥣 Misión 3: Sopa de Letras</div>')
        body.append('<p class="inst">Encuentra y encierra las siguientes palabras.</p>')
        chips = "".join(f'<span class="word-chip">{p}</span>' for p in palabras_sopa)
        body.append(f'<div class="word-bank">{chips}</div>')
        body.append(_sopa_html(palabras_sopa))
        body.append('</div>')

    # ── Misión 4: El Pergamino (fill-in) ────────────────────────────────────
    per = misiones.get("pergamino", {})
    frase = str(per.get("frase_con_espacios","")).strip()
    palabras_secretas = per.get("palabras_secretas", [])
    if frase:
        body.append('<div class="section">')
        body.append('<div class="section-header purple">📜 Misión 4: El Pergamino</div>')
        body.append('<p class="inst">Completa la frase usando las palabras del banco.</p>')
        # Render blanks as styled underlines
        frase_html = re.sub(r'_{3,}', '<span class="fill-blank">&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;&nbsp;</span>', frase)
        body.append(f'<div class="fill-box">{frase_html}</div>')
        if palabras_secretas:
            chips = "".join(f'<span class="word-chip">{p}</span>' for p in palabras_secretas)
            body.append(f'<p class="inst">Banco de palabras:</p><div class="word-bank">{chips}</div>')
        body.append('</div>')

    # ── Misión 5: El Lienzo (drawing) ───────────────────────────────────────
    lienzo_inst = str(misiones.get("lienzo","Crea tu obra de arte basándote en lo aprendido.")).replace("**","")
    body.append('<div class="section">')
    body.append('<div class="section-header orange">🎨 Misión 5: El Lienzo</div>')
    body.append(f'<p>{lienzo_inst}</p>')
    body.append('<div class="draw-box"><span class="draw-label">✏️ Dibuja o escribe tu respuesta aquí</span></div>')
    body.append('</div>')

    # ── Teacher adaptations (separate page) ─────────────────────────────────
    adapt = ficha.get("adaptaciones_por_mision", [])
    if adapt:
        body.append('<div class="teacher-page">')
        body.append('<div class="teacher-banner">📋 ANEXO DOCENTE — Adaptaciones por Misión</div>')
        body.append('<p class="confidential">Esta página es solo para el docente. No imprimir para los estudiantes.</p>')
        if isinstance(adapt[0], dict):
            cols = list(adapt[0].keys())
            body.append('<table class="match-table"><tr>')
            for c in cols:
                body.append(f'<th>{c}</th>')
            body.append('</tr>')
            for row in adapt:
                body.append('<tr>')
                for c in cols:
                    body.append(f'<td>{row.get(c,"")}</td>')
                body.append('</tr>')
            body.append('</table>')
        else:
            body.append('<ul>' + "".join(f"<li>{str(a)}</li>" for a in adapt) + '</ul>')
        body.append('</div>')

    return _render_html_doc(
        title   = f"Ficha Gamificada — {tema}",
        materia = "Ficha Gamificada",
        tema    = tema,
        body_html = "\n".join(body),
        theme   = theme,
    )


# ─────────────────────────────────────────────────────────────────────────────
# HTML PLAN DOCUMENT GENERATORS  (teachers / admin — not student worksheets)
# ─────────────────────────────────────────────────────────────────────────────
_PLAN_HTML_CSS = """
*{box-sizing:border-box;margin:0;padding:0}
body{font-family:'Segoe UI',Calibri,Arial,sans-serif;font-size:11pt;color:#1A1A1A;background:#fff;padding:28px 36px}
.doc-header{display:flex;align-items:center;gap:16px;background:#1A3A6A;color:#fff;padding:14px 20px;border-radius:8px;margin-bottom:8px;-webkit-print-color-adjust:exact;print-color-adjust:exact}
.doc-header-main{font-size:14pt;font-weight:700}
.doc-header-sub{font-size:9pt;color:#C8E6C9;margin-top:3px}
.doc-title{font-size:18pt;font-weight:700;color:#CC2200;text-align:center;margin:14px 0 6px;text-decoration:underline;text-underline-offset:4px;text-transform:uppercase}
h1{color:#1A3A6A;font-size:14pt;margin:18px 0 6px;border-bottom:2px solid #E3F0FF;padding-bottom:4px}
h2{color:#1A3A6A;font-size:12pt;margin:12px 0 5px}
h3{color:#2E7D32;font-size:11pt;margin:9px 0 4px}
p{margin:5px 0 8px;line-height:1.5}
ul,ol{margin:6px 0 10px;padding-left:24px}
li{margin:3px 0;line-height:1.5}
.banner{background:#1A3A6A;color:#fff;padding:7px 14px;border-radius:6px;font-weight:700;font-size:11pt;margin:14px 0 8px;-webkit-print-color-adjust:exact;print-color-adjust:exact}
.banner.teal{background:#007A6E}
.banner.navy2{background:#1A5276}
table{border-collapse:collapse;width:100%;margin:10px 0}
th{background:#1A3A6A;color:#fff;padding:7px 10px;text-align:left;font-size:10pt;-webkit-print-color-adjust:exact;print-color-adjust:exact}
td{border:1px solid #DDD;padding:7px 10px;font-size:10pt}
tr:nth-child(even) td{background:#F1F8E9}
blockquote{border-left:4px solid #007A6E;padding:6px 12px;background:#F0FAF8;font-style:italic;margin:8px 0}
hr{border:none;border-top:1px solid #DDD;margin:16px 0}
@media print{@page{size:A4;margin:15mm 12mm}body{padding:0;font-size:10pt}}
"""

def _render_plan_doc(title: str, subtitle: str, body_html: str) -> bytes:
    fecha = datetime.now().strftime("%d/%m/%Y")
    logo_tag = ""
    if os.path.exists(LOGO_PATH):
        import base64 as _b64
        with open(LOGO_PATH, "rb") as _f:
            _b64s = _b64.b64encode(_f.read()).decode()
        logo_tag = f'<img src="data:image/png;base64,{_b64s}" style="height:44px;width:auto">'
    html = f"""<!DOCTYPE html>
<html lang="es"><head><meta charset="UTF-8"><title>{title}</title>
<style>{_PLAN_HTML_CSS}</style></head><body>
<div class="doc-header">{logo_tag}
  <div style="flex:1">
    <div class="doc-header-main">Las Palmas School — PRIA</div>
    <div class="doc-header-sub">{subtitle}</div>
  </div>
  <div style="color:#C8E6C9;font-size:9pt">{fecha}</div>
</div>
<div class="doc-title">{title}</div>
{body_html}
</body></html>"""
    return html.encode("utf-8")


def _md_to_plan_html(texto: str) -> str:
    if not texto:
        return ""
    out = []; in_ul = False
    for line in texto.split("\n"):
        s = line.strip()
        if _es_autopresentacion(s): continue
        if not s:
            if in_ul: out.append("</ul>"); in_ul = False
            continue
        if s in ("---", "___", "***") or (len(s) >= 3 and all(c == "-" for c in s)):
            if in_ul: out.append("</ul>"); in_ul = False
            out.append("<hr>"); continue
        s_h = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', s)
        s_h = re.sub(r'\*(.+?)\*',     r'<em>\1</em>', s_h)
        if   s.startswith("### "): out.append(f"<h3>{s_h[4:]}</h3>") if not in_ul else (out.append("</ul>"), out.append(f"<h3>{s_h[4:]}</h3>")); in_ul = False
        elif s.startswith("## "):  out.append(f"<h2>{s_h[3:]}</h2>") if not in_ul else (out.append("</ul>"), out.append(f"<h2>{s_h[3:]}</h2>")); in_ul = False
        elif s.startswith("# "):   out.append(f"<h1>{s_h[2:]}</h1>") if not in_ul else (out.append("</ul>"), out.append(f"<h1>{s_h[2:]}</h1>")); in_ul = False
        elif s.startswith("> "):
            if in_ul: out.append("</ul>"); in_ul = False
            out.append(f"<blockquote>{s_h[2:]}</blockquote>")
        elif s.startswith("- ") or s.startswith("* "):
            if not in_ul: out.append("<ul>"); in_ul = True
            out.append(f"<li>{s_h[2:]}</li>")
        else:
            if in_ul: out.append("</ul>"); in_ul = False
            out.append(f"<p>{s_h}</p>")
    if in_ul: out.append("</ul>")
    return "\n".join(out)


def _lst(v) -> list:
    return v if isinstance(v, list) else ([v] if v else [])


def generar_html_sintesis(res_m0a: dict, diagnosticos: str) -> bytes:
    unidad = res_m0a.get("unidad_sintetizada", {})
    titulo = unidad.get("titulo", "Síntesis de Unidad")
    body   = f"<h1>Perfil Neuro-Inclusivo del Aula</h1><p>{diagnosticos}</p>"
    body  += "<h1>Temas de la Unidad</h1>"
    for tema in unidad.get("temas_desarrollados", []):
        nombre = tema.get("nombre", "Tema")
        body  += f"<h2>{nombre}</h2>"
        ccs = _lst(tema.get("conceptos_clave") or tema.get("conceptos") or [])
        if ccs:
            body += "<strong>Conceptos Clave:</strong><ul>" + "".join(f"<li>{c}</li>" for c in ccs) + "</ul>"
        ims = _lst(tema.get("inteligencias_sugeridas") or tema.get("inteligencias") or [])
        if ims:
            body += "<strong>Inteligencias Múltiples:</strong><ul>" + "".join(f"<li>{i}</li>" for i in ims) + "</ul>"
    body += f'<h1>Notas DUA del Docente</h1>{_md_to_plan_html(unidad.get("notas_docente",""))}'
    body += f'<h1>Proyecto ABP — Resumen</h1>{_md_to_plan_html(unidad.get("proyecto_pbl",""))}'
    return _render_plan_doc("Síntesis Curricular", titulo, body)


def generar_html_abp(res_m0b: str, res_m0c: str, titulo_unidad: str) -> bytes:
    body  = '<div class="banner">📋 Diseño del Proyecto ABP</div>'
    body += _md_to_plan_html(res_m0b or "No generado.")
    body += '<div class="banner teal">📊 Rúbrica y Fichas de Proceso</div>'
    body += _md_to_plan_html(res_m0c or "No generado.")
    return _render_plan_doc("Proyecto ABP", titulo_unidad, body)


def generar_html_plan_clase(res_m1a: dict, tema: str, diagnosticos: str) -> bytes:
    body = ""
    verbos = res_m1a.get("mapa_cognitivo", {}).get("verbos", [])
    if verbos:
        body += "<h1>Verbos Operativos (Taxonomía Bloom)</h1><ul>"
        body += "".join(f"<li>{v}</li>" for v in verbos) + "</ul>"
    im = res_m1a.get("inteligencias_multiples", [])
    if im and isinstance(im, list) and im and isinstance(im[0], dict):
        cols  = list(im[0].keys())
        body += "<h1>Inteligencias Múltiples</h1><table>"
        body += "<tr>" + "".join(f"<th>{c}</th>" for c in cols) + "</tr>"
        for row in im:
            body += "<tr>" + "".join(f"<td>{row.get(c,'')}</td>" for c in cols) + "</tr>"
        body += "</table>"
    body += '<div class="banner">⏱️ Secuencia Didáctica</div>'
    for bloque in res_m1a.get("secuencia_didactica", {}).get("bloques", []):
        nb = bloque.get("nombre", "Bloque"); dur = bloque.get("duracion", "?")
        body += f"<h2>{nb} — {dur} min</h2><p>{bloque.get('objetivo','')}</p>"
        if "nota" in bloque:
            body += f"<blockquote>💡 {bloque['nota']}</blockquote>"
    body += '<div class="banner teal">🧠 Directrices DUA / Neuro-Inclusión</div><ul>'
    body += "".join(f"<li>{d}</li>" for d in _lst(res_m1a.get("dua_neuroinclusion", []))) + "</ul>"
    adapt = res_m1a.get("tabla_adaptaciones_clase", [])
    if adapt and isinstance(adapt, list) and adapt and isinstance(adapt[0], dict):
        cols  = list(adapt[0].keys())
        body += '<div class="banner navy2">♿ Adaptaciones por Diagnóstico</div>'
        body += "<table><tr>" + "".join(f"<th>{c}</th>" for c in cols) + "</tr>"
        for row in adapt:
            body += "<tr>" + "".join(f"<td>{row.get(c,'')}</td>" for c in cols) + "</tr>"
        body += "</table>"
    body += f"<h1>Perfil del Aula</h1><p>{res_m1a.get('perfil_aula_resumido', diagnosticos)}</p>"
    return _render_plan_doc("Plan de Clase · 45 minutos", tema, body)


def generar_html_pdc(resultado_pdc: str, teacher: str, subject: str, level: str) -> bytes:
    return _render_plan_doc("PDC Trimestral", f"{subject} — {teacher} — {level}",
                            _md_to_plan_html(resultado_pdc or "No generado."))


# ─────────────────────────────────────────────────────────────────────────────
# PANEL STREAMLIT
# ─────────────────────────────────────────────────────────────────────────────
def render_panel_exportacion(ss, diagnosticos: str):
    import streamlit as st
    try:
        from slide_generator import generar_pptx_con_claude, node_disponible, instalar_pptxgenjs
    except ImportError as _e:
        # Graceful fallback — premium engine unavailable
        def node_disponible(): return False
        def instalar_pptxgenjs(): return str(_e)
        generar_pptx_con_claude = None

    unidad_titulo = ""
    if ss.res_m0a and "unidad_sintetizada" in ss.res_m0a:
        unidad_titulo = ss.res_m0a["unidad_sintetizada"].get("titulo", "Unidad")
    tema_activo = ss.get("tema_activo", "tema")
    fecha_str   = datetime.now().strftime("%Y%m%d")
    HTML = "text/html"
    PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    anthropic_key = st.secrets.get("ANTHROPIC_API_KEY", "")
    node_ok       = node_disponible()

    st.markdown("### 📤 Exportar — Plantilla Las Palmas School")
    st.caption("Logo oficial. Listo para imprimir. Solo se habilitan los materiales ya generados.")

    if not node_ok or not anthropic_key:
        with st.expander("⚙️ Configuración del Motor de Diapositivas", expanded=False):
            if not anthropic_key:
                st.warning("⚠️ Falta `ANTHROPIC_API_KEY` en `secrets.toml`.")
            if not node_ok:
                st.warning("⚠️ Node.js o pptxgenjs no encontrado.")
                if st.button("📦 Instalar pptxgenjs ahora"):
                    with st.spinner("Instalando…"):
                        st.info(instalar_pptxgenjs())
                else:
                    st.caption("Ejecuta: `npm install -g pptxgenjs`")
            st.info("Diapositivas en modo básico (python-pptx) hasta configurar el motor premium.")

    st.markdown("#### 🏗️ Plan de Unidad y ABP")
    c1, c2, _ = st.columns(3)
    with c1:
        if ss.res_m0a:
            st.download_button("📄 Síntesis Curricular (.html)",
                data=generar_html_sintesis(ss.res_m0a, diagnosticos),
                file_name=f"LasP_Sintesis_{fecha_str}.html", mime=HTML,
                use_container_width=True)
        else:
            st.button("📄 Síntesis Curricular (.html)", disabled=True,
                      use_container_width=True, help="Genera la Síntesis primero.")
    with c2:
        if ss.res_m0b:
            st.download_button("📋 Proyecto ABP (.html)",
                data=generar_html_abp(ss.res_m0b, ss.res_m0c or "", unidad_titulo),
                file_name=f"LasP_ABP_{fecha_str}.html", mime=HTML,
                use_container_width=True)
        else:
            st.button("📋 Proyecto ABP (.html)", disabled=True,
                      use_container_width=True, help="Diseña el Proyecto ABP primero.")

    # ── Selector de estilo de plantilla HTML ────────────────────────────────
    with st.expander("🎨 Estilo de plantilla (fichas y evaluaciones HTML)", expanded=False):
        _theme_options = list(WORKSHEET_THEMES.keys())
        _theme_labels  = [WORKSHEET_THEMES[k]["label"] for k in _theme_options]
        _theme_idx = st.radio(
            "Elige el estilo visual:",
            options=range(len(_theme_options)),
            format_func=lambda i: _theme_labels[i],
            key="export_theme_idx",
            horizontal=False,
        )
        html_theme = _theme_options[_theme_idx]
        st.caption(f"Vista previa activa: **{_theme_labels[_theme_idx]}**")

    st.markdown("#### 🚀 Plan de Clase y Fichas")
    c4, c5, c6 = st.columns(3)
    with c4:
        if ss.res_m1a:
            st.download_button("📄 Plan de Clase (.html)",
                data=generar_html_plan_clase(ss.res_m1a, tema_activo, diagnosticos),
                file_name=f"LasP_PlanClase_{fecha_str}.html", mime=HTML,
                use_container_width=True)
        else:
            st.button("📄 Plan de Clase (.html)", disabled=True,
                      use_container_width=True, help="Genera el Plan de Clase primero.")
    with c5:
        if ss.res_m1b:
            if node_ok and anthropic_key:
                if st.button("✨ Generar Diapositivas (.pptx)",
                             use_container_width=True, type="primary"):
                    with st.spinner("🎨 Claude diseñando las diapositivas… (30–60 seg)"):
                        try:
                            slides = (ss.res_m1b if isinstance(ss.res_m1b, list)
                                      else ss.res_m1b.get("diapositivas")
                                         or ss.res_m1b.get("slides")
                                         or list(ss.res_m1b.values())[0])
                            pptx_bytes = generar_pptx_con_claude(
                                slides_data=slides, tema=tema_activo,
                                unidad=unidad_titulo, anthropic_api_key=anthropic_key)
                            st.session_state["_pptx_cache"] = pptx_bytes
                            st.success("✅ ¡Diapositivas generadas!")
                        except Exception as e:
                            st.error(f"❌ {str(e)[:400]}")
                if st.session_state.get("_pptx_cache"):
                    st.download_button("⬇️ Descargar Diapositivas (.pptx)",
                        data=st.session_state["_pptx_cache"],
                        file_name=f"LasP_Slides_{fecha_str}.pptx", mime=PPTX,
                        use_container_width=True)
            else:
                st.download_button("🖼️ Diapositivas (.pptx) — modo básico",
                    data=generar_pptx_diapositivas(ss.res_m1b, tema_activo, unidad_titulo),
                    file_name=f"LasP_Slides_{fecha_str}.pptx", mime=PPTX,
                    use_container_width=True)
        else:
            st.button("🖼️ Diapositivas (.pptx)", disabled=True,
                      use_container_width=True, help="Genera las Diapositivas primero.")
    with c6:
        if ss.res_m1c:
            st.download_button("🎮 Ficha Gamificada (.html)",
                data=generar_html_ficha(ss.res_m1c, tema_activo, theme=html_theme).encode("utf-8"),
                file_name=f"LasP_Ficha_{fecha_str}.html", mime=HTML,
                use_container_width=True)
        else:
            st.button("🎮 Ficha Gamificada (.html)", disabled=True,
                      use_container_width=True, help="Genera la Ficha primero.")

    st.markdown("#### 📝 Evaluaciones")
    c7, _, __ = st.columns(3)
    with c7:
        if ss.res_m2a:
            st.download_button("📝 Pop Quiz + Guía Tutor (.html)",
                data=generar_html_evaluaciones(ss.res_m2a, ss.res_m2b or "", tema_activo, theme=html_theme).encode("utf-8"),
                file_name=f"LasP_Evaluaciones_{fecha_str}.html", mime=HTML,
                use_container_width=True)
        else:
            st.button("📝 Pop Quiz + Guía Tutor (.html)", disabled=True,
                      use_container_width=True, help="Genera el Pop Quiz primero.")
